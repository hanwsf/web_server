#!/usr/bin/env python3
"""
语义分析调度系统
- 扫描和注册 Skills
- 基于LLM的意图分析
- 调度 Skills/Agents/Subagents
todo改进: 文件夹固定，模型固定为minimax-m2.1
"""

import os
import re
import concurrent.futures
from urllib.parse import quote as url_quote
import yaml
import json
# import logging
from loguru import logger
from pathlib import Path
from typing import Dict, List, Optional, Any
from dataclasses import dataclass, field
from datetime import datetime

# logger = logging.getLogger(__name__)

# ============ 模型映射 ============
# 将前端模型选择值映射到实际的模型标识符
MODEL_MAPPING = {
    # Claude 模型系列
    "opus": "claude-opus-4-5-20251101",
    # "sonnet": "sonnet",
    "haiku": "claude-haiku-4-5-20251001",

    # GLM 模型系列
    "glm4.7": "nvidia/z-ai/glm4.7",

    # Deepseek 模型系列
    "deepseek-v3.1-terminus": "nvidia/deepseek-ai/deepseek-v3.1-terminus",
    "deepseek-v3.2": "nvidia/deepseek-ai/deepseek-v3.2",

    # MiniMax 模型系列
    "minimax-m2.1": "nvidia/minimaxai/minimax-m2.1",

    # Deepseek Chat
    "deepseek-chat": "deepseek/deepseek-chat",
}

# ============ 知识库路径 ============
WORK_DIR = Path(os.environ.get('WORK_DIR', '/home/will/Downloads/opencode_p'))
KB_PATHS = {
    "KB": str(WORK_DIR / "KB"),
    "KBGX": str(WORK_DIR / "KBGX"),
    "KBW": str(WORK_DIR / "KBW")
}


def normalize_model(model: str) -> str:
    """将前端选择的模型值标准化为支持的格式"""
    if not model:
        return "haiku"  # 使用 Claude Haiku 作为默认模型（快速稳定）

    # 如果有映射，使用映射值
    if model in MODEL_MAPPING:
        normalized = MODEL_MAPPING[model]
        logger.info(f"[MODEL] 将前端模型 '{model}' 映射到 '{normalized}'")
        return normalized

    # 否则直接使用（假设前端已经传入了正确的格式）
    logger.info(f"[MODEL] 使用前端模型 '{model}'")
    return model

# ============ 数据结构 ============

@dataclass
class SkillOperation:
    """Skill的操作定义"""
    name: str
    semantic_triggers: List[str]
    script: str
    params: List[str]

@dataclass
class SkillInfo:
    """Skill信息"""
    name: str
    description: str
    path: Path
    operations: Dict[str, SkillOperation] = field(default_factory=dict)
    keywords: List[str] = field(default_factory=list)
    use_cases: List[str] = field(default_factory=list)

@dataclass
class AgentInfo:
    """Agent信息"""
    name: str
    description: str
    subagent_type: str
    use_cases: List[str] = field(default_factory=list)

@dataclass
class Intent:
    """意图分析结果"""
    intent_type: str  # skill, agent, subagent, direct
    target: str  # skill name / agent type
    operation: Optional[str] = None  # specific operation for skills
    params: Dict[str, Any] = field(default_factory=dict)
    confidence: float = 1.0
    reasoning: str = ""
    is_complex: bool = False  # 是否为复杂多步骤任务
    subtasks: List['SubTask'] = field(default_factory=list)  # 分解后的子任务

@dataclass
class SubTask:
    """子任务定义"""
    task_id: str  # 唯一ID，格式: step_1, step_2等
    name: str  # 任务名称
    intent_type: str  # skill/agent/subagent/direct
    target: str  # 目标skill/agent
    operation: Optional[str] = None
    params: Dict[str, Any] = field(default_factory=dict)
    timeout_seconds: int = 300  # 5分钟默认超时
    depends_on: List[str] = field(default_factory=list)  # 依赖的子任务ID列表
    status: str = "pending"  # pending/running/completed/failed/timeout
    result: Optional[Any] = None
    error: Optional[str] = None
    start_time: Optional[datetime] = None
    end_time: Optional[datetime] = None

@dataclass
class ComplexRequestMetadata:
    """复杂请求的元数据"""
    is_complex: bool = False
    task_count: int = 0
    total_timeout: int = 0  # 总超时时间（秒）
    subtasks_list: List[SubTask] = field(default_factory=list)
    execution_dag: Dict[str, List[str]] = field(default_factory=dict)  # 依赖图

# ============ 任务分解器 ============

class TaskDecomposer:
    """复杂请求分解器 - 将复杂请求分解为多个子任务"""

    def __init__(self, model_caller, discovery: Optional['SkillDiscovery'] = None):
        self.model_caller = model_caller
        self.discovery = discovery
        self.complex_patterns = self._init_complex_patterns()

    def _init_complex_patterns(self) -> Dict:
        """初始化复杂请求的特征模式"""
        return {
            "connectors": ["和", "然后", "之后", "完成后", "最后", "接着", "随后", "+", "同时"],
            "multi_tool_indicators": ["利用", "使用", "调用", "转换", "发送", "生成"],
            "format_conversion": ["转为", "转成", "格式", "转换", "输出"],
            "distribution": ["发送", "邮件", "发到", "发给", "email", "上传", "下载"],
        }

    def is_complex_request(self, user_input: str) -> bool:
        """检测请求是否为复杂多步骤任务"""
        patterns = self.complex_patterns

        # 检查连接词
        has_connector = any(conn in user_input for conn in patterns["connectors"])
        # 检查多工具指示符
        has_multi_tools = user_input.count("skill") > 1 or \
                          any(tool in user_input for tool in ["知识库", "邮件", "Word", "docx", "pptx"])

        return has_connector or (has_multi_tools and len(user_input) > 50)
#复杂
    def decompose(self, user_input: str, model: str = "haiku") -> ComplexRequestMetadata:
        """分解复杂请求为子任务"""
        logger.info(f"[Decomposer] 开始分解请求: {user_input[:60]}...")

        if not self.is_complex_request(user_input):
            logger.info(f"[Decomposer] 非复杂请求，跳过分解")
            return ComplexRequestMetadata(is_complex=False)

        # 使用LLM进行请求分解
        subtasks = self._llm_decompose(user_input, model)

        if not subtasks:
            logger.warning(f"[Decomposer] LLM分解失败，尝试本地规则")
            # subtasks = self._local_decompose(user_input)

        # 构建执行DAG和元数据
        metadata = self._build_execution_plan(subtasks, user_input)
        logger.info(f"[Decomposer] 分解完成: {len(metadata.subtasks_list)} 个子任务:{str(metadata.subtasks_list)}")

        return metadata

    def _llm_decompose(self, user_input: str, model: str) -> Optional[List[SubTask]]:
        """使用LLM分解复杂请求"""
        # 构建可用资源列表
        available_resources = self._get_available_resources()

        system_prompt = f"""你是一个任务分解专家。将用户的复杂多步骤请求分解为一系列清晰的子任务。

重要约束：
1. 只能使用下面列出的 Skills和Agents
2. 如果没有合适的 Skill/Agent，使用 type="direct"
3. 不要创建虚拟的 Agent/Skill，否则会导致任务失败

【type="direct" 时的 target 规则 - 必须严格遵守】
- 网络搜索 → target="WebSearch_dk"
- 发送邮件 → target="send_mail"
- 从GitHub获取代码/仓库内容 → target="WebFetch" (operation填写完整URL)
- 其他由Claude直接处理的任务 → target="claude"

【可用的 Skills】
{available_resources['skills_list']}

【可用的 Agents】
{available_resources['agents_list']}

每个子任务应该包含：
1. 任务名称：简洁描述
2. 任务类型：skill/agent/direct
3. 目标：具体要调用的skill或agent，或 direct 类型的特定 target（见上方规则）
4. 依赖：此任务依赖哪些前序任务（用任务ID引用，如step_1, step_2）

返回格式为JSON数组，每项包含：
{{
  "name": "任务名称",
  "type": "skill|agent|direct",
  "target": "skill名/agent名/WebSearch_dk/send_mail/WebFetch/claude",
  "operation": "操作描述或URL",
  "depends_on": ["step_1"]
}}

注意：如果 type 是 "skill"，target 必须是上面列出的 Skills 之一；如果 type 是 "agent"，target 必须是上面列出的 Agents 之一。
【重要】: 每个子任务必须是独立的单一操作，如：                                                                                                                                                
- "获取GitHub内容" 和 "创作文章" 应该是两个独立的步骤                                                                                                                                     
- "创建Word" 和 "发送邮件" 应该是两个独立的步骤                                                                                                                                           
- 不要合并多个操作到一个步骤中  """

        user_prompt = f"""请分解以下请求为子任务：

{user_input}

请严格遵守约束条件，只使用可用的 Skills/Agents。
请以JSON数组形式返回，只返回JSON，不要其他内容。"""

        try:
            response = self.model_caller(system_prompt, user_prompt, model=model, temperature=0.1)
            logger.debug(f"[Decomposer] LLM响应: {response}...")

            result = self._extract_json_array(response)
            logger.debug(f"[Decomposer] LLM响应result: {str(result)}...")
            if result:
                subtasks = []
                for i, task_data in enumerate(result, 1):
                    subtask = SubTask(
                        task_id=f"step_{i}",
                        name=task_data.get("name", f"Task {i}"),
                        intent_type=task_data.get("type", "direct"),
                        target=task_data.get("target", "claude"),
                        depends_on=task_data.get("depends_on", []),
                        timeout_seconds=300
                    )
                    subtasks.append(subtask)
                logger.info(f"[Decomposer] LLM分解结果: {len(subtasks)} 个子任务")
                return subtasks
        except Exception as e:
            logger.warning(f"[Decomposer] LLM分解异常: {e}")

        return None

    def _get_available_resources(self) -> Dict:
        """获取可用的 Skills 和 Agents 列表"""
        skills_list = "无"
        agents_list = "无"

        if self.discovery:
            # 获取 Skills 列表
            if self.discovery.skills:
                skills_list = "\n".join([f"- {name}: {skill.description[:60]}"
                                        for name, skill in self.discovery.skills.items()])

            # 获取 Agents 列表
            if self.discovery.agents:
                agents_list = "\n".join([f"- {name}: {agent.description}"
                                        for name, agent in self.discovery.agents.items()])

        return {
            "skills_list": skills_list,
            "agents_list": agents_list
        }
#llm失败后的本地规则，
    # def _local_decompose(self, user_input: str) -> List[SubTask]:
    #     """本地规则分解（后备方案）"""
    #     subtasks = []
    #     input_lower = user_input.lower()

    #     # 检测微信文章+Word+邮件的模式
    #     if "wechat" in input_lower and ("word" in input_lower or "docx" in input_lower) and "邮件" in input_lower:
    #         subtasks = [
    #             SubTask(task_id="step_1", name="获取GitHub内容", intent_type="direct",
    #                    target="claude", timeout_seconds=300),
    #             SubTask(task_id="step_2", name="生成微信文章", intent_type="skill",
    #                    target="wechat-article-editor", depends_on=["step_1"], timeout_seconds=600),
    #             SubTask(task_id="step_3", name="转换为Word格式", intent_type="skill",
    #                    target="docx", depends_on=["step_2"], timeout_seconds=600),
    #             SubTask(task_id="step_4", name="发送邮件", intent_type="direct",
    #                    target="claude", depends_on=["step_3"], timeout_seconds=300),
    #         ]

    #     return subtasks

    def _extract_json_array(self, text: str) -> Optional[List[Dict]]:
        """从文本中提取JSON数组"""
        try:
            # 查找JSON数组
            match = re.search(r'\[\s*{.*?}\s*\]', text, re.DOTALL)
            if match:
                json_str = match.group(0)
                return json.loads(json_str)
        except Exception as e:
            logger.debug(f"[Decomposer] JSON提取失败: {e}")
        return None

    def _build_execution_plan(self, subtasks: List[SubTask], user_input: str) -> ComplexRequestMetadata:
        """构建执行计划和DAG"""
        # 计算总超时
        total_timeout = len(subtasks) * 480

        # 构建依赖图
        dag = {}
        for task in subtasks:
            dag[task.task_id] = task.depends_on

        metadata = ComplexRequestMetadata(
            is_complex=True,
            task_count=len(subtasks),
            total_timeout=total_timeout,
            subtasks_list=subtasks,
            execution_dag=dag
        )
        logger.info(f"metadata:{str(metadata)}")
        return metadata

# ============ 任务执行协调器 ============

class TaskExecutionOrchestrator:
    """任务执行协调器 - 管理多个子任务的并发执行、超时和依赖"""

    def __init__(self):
        self.running_tasks: Dict[str, SubTask] = {}
        self.completed_tasks: Dict[str, SubTask] = {}

    async def execute_plan(self, metadata: ComplexRequestMetadata,
                          executor_func) -> ComplexRequestMetadata:
        """执行任务计划"""
        logger.info(f"[Orchestrator] 开始执行任务计划: {metadata.task_count} 个子任务")

        with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
            futures = {}

            # 拓扑排序执行
            executed = set()
            while len(executed) < len(metadata.subtasks_list):
                # 找到可以执行的任务（所有依赖都已完成）
                ready_tasks = [
                    task for task in metadata.subtasks_list
                    if task.task_id not in executed and
                    all(dep in executed for dep in task.depends_on)
                ]

                if not ready_tasks:
                    # 如果没有就绪任务但还有未执行的，说明有循环依赖
                    logger.error("[Orchestrator] 检测到循环依赖或异常")
                    break

                # 提交就绪任务
                for task in ready_tasks:
                    future = executor.submit(
                        self._execute_subtask,
                        task,
                        executor_func
                    )
                    futures[task.task_id] = (task, future)

                # 等待至少一个任务完成
                for task_id, (task, future) in list(futures.items()):
                    if task_id not in executed:
                        try:
                            future.result(timeout=task.timeout_seconds)
                            executed.add(task_id)
                            logger.info(f"[Orchestrator] 任务 {task_id} 完成")
                        except concurrent.futures.TimeoutError:
                            task.status = "timeout"
                            task.error = f"超时: {task.timeout_seconds}秒"
                            executed.add(task_id)
                            logger.error(f"[Orchestrator] 任务 {task_id} 超时")
                        except Exception as e:
                            task.status = "failed"
                            task.error = str(e)
                            executed.add(task_id)
                            logger.error(f"[Orchestrator] 任务 {task_id} 失败: {e}")

        return metadata

    def _execute_subtask(self, task: SubTask, executor_func) -> Any:
        """执行单个子任务"""
        task.start_time = datetime.now()
        task.status = "running"

        try:
            logger.info(f"[Orchestrator] 执行子任务: {task.task_id} ({task.name})")
            result = executor_func(task)

            task.status = "completed"
            task.result = result
            task.end_time = datetime.now()

            elapsed = (task.end_time - task.start_time).total_seconds()
            logger.info(f"[Orchestrator] 任务 {task.task_id} 完成 (耗时: {elapsed:.1f}秒)")

            return result
        except Exception as e:
            task.status = "failed"
            task.error = str(e)
            task.end_time = datetime.now()
            logger.error(f"[Orchestrator] 任务 {task.task_id} 执行异常: {e}")
            raise

# ============ Skill 发现器 ============

class SkillDiscovery:
    """发现和注册 Skills"""

    def __init__(self, skills_dir: Path):
        self.skills_dir = Path(skills_dir)
        self.skills: Dict[str, SkillInfo] = {}
        self.agents: Dict[str, AgentInfo] = {}
        self._initialize_builtin_agents()

    def _initialize_builtin_agents(self):
        """初始化内置Agent定义"""
        self.agents = {
            "Bash": AgentInfo(
                name="Bash",
                description="命令执行专家 - 运行shell命令、git操作、文件操作",
                subagent_type="Bash",
                use_cases=[
                    "运行命令行工具",
                    "git操作（clone, commit, push等）",
                    "文件系统操作（移动, 复制, 删除）",
                    "安装依赖、构建项目"
                ]
            ),
            "general-purpose": AgentInfo(
                name="general-purpose",
                description="通用研究代理 - 复杂搜索、代码探索、多步骤任务",
                subagent_type="general-purpose",
                use_cases=[
                    "需要多轮搜索的复杂问题",
                    "跨多个文件或模块的代码分析",
                    "需要综合多个信息源的任务"
                ]
            ),
            "Explore": AgentInfo(
                name="Explore",
                description="代码库探索代理 - 快速浏览代码结构、搜索特定功能",
                subagent_type="Explore",
                use_cases=[
                    "查找文件或函数定义",
                    "理解项目结构",
                    "搜索代码中的特定模式",
                    "快速查看某个功能在代码库中的位置"
                ]
            ),
            "Plan": AgentInfo(
                name="Plan",
                description="架构规划代理 - 设计复杂功能的实现计划",
                subagent_type="Plan",
                use_cases=[
                    "设计新功能的架构",
                    "规划复杂的重构任务",
                    "需要多种实现方案对比的任务"
                ]
            ),
            "claude-code-guide": AgentInfo(
                name="claude-code-guide",
                description="Claude使用指南 - 解答关于Claude使用的问题",
                subagent_type="claude-code-guide",
                use_cases=[
                    "询问Claude的功能",
                    "了解如何使用特定工具",
                    "获取使用技巧和最佳实践"
                ]
            )
        }

    def discover_skills(self) -> Dict[str, SkillInfo]:
        """发现所有Skills"""
        logger.info(f"[Discovery] 开始扫描 {self.skills_dir}")

        for skill_dir in self.skills_dir.iterdir():
            if not skill_dir.is_dir() or skill_dir.name.startswith('.'):
                continue

            skill_md = skill_dir / "SKILL.md"
            if not skill_md.exists():
                continue

            skill_info = self._parse_skill_md(skill_md)
            if skill_info:
                self.skills[skill_info.name] = skill_info
                logger.info(f"[Discovery] 发现 Skill: {skill_info.name}")

        logger.info(f"[Discovery] 共发现 {len(self.skills)} 个 Skills")
        return self.skills

    def _parse_skill_md(self, md_path: Path) -> Optional[SkillInfo]:
        """解析 SKILL.md 文件"""
        try:
            content = md_path.read_text(encoding='utf-8')

            # 提取 YAML frontmatter
            yaml_match = re.match(r'^---\n(.*?)\n---', content, re.DOTALL)
            if not yaml_match:
                # 没有 YAML frontmatter，使用简单的 name 提取
                name = md_path.parent.name
                description = content.split('\n\n')[0][:300] if content else ""
                return SkillInfo(name=name, description=description, path=md_path.parent)

            yaml_content = yaml_match.group(1)
            metadata = yaml.safe_load(yaml_content)

            name = metadata.get('name', md_path.parent.name)
            description = metadata.get('description', '')
            keywords = []
            use_cases = []
            operations = {}

            # 提取 keywords
            if 'keywords' in metadata:
                keywords = [k.strip() for k in metadata['keywords'] if k.strip()]

            # 提取 use_cases (从description中分析)
            use_cases = self._extract_use_cases(description)

            # 提取 operations (如果有 operations 定义)
            if 'operations' in metadata:
                for op_name, op_data in metadata['operations'].items():
                    operations[op_name] = SkillOperation(
                        name=op_name,
                        semantic_triggers=op_data.get('semantic_triggers', []),
                        script=op_data.get('script', ''),
                        params=op_data.get('params', [])
                    )

            return SkillInfo(
                name=name,
                description=description,
                path=md_path.parent,
                operations=operations,
                keywords=keywords,
                use_cases=use_cases
            )

        except Exception as e:
            logger.warning(f"[Discovery] 解析失败 {md_path}: {e}")
            return None

    def _extract_use_cases(self, description: str) -> List[str]:
        """从描述中提取使用场景"""
        use_cases = []
        # 匹配数字列表项
        pattern = r'[（(]\s*\d+\s*[)）]'
        if re.search(pattern, description):
            parts = re.split(pattern, description)
            for part in parts[1:]:  # 跳过第一个非编号部分
                if part.strip():
                    use_cases.append(part.strip()[:100])
        return use_cases

    def get_skill_by_name(self, name: str) -> Optional[SkillInfo]:
        """按名称获取Skill"""
        return self.skills.get(name)

    def get_agent_by_name(self, name: str) -> Optional[AgentInfo]:
        """按名称获取Agent"""
        return self.agents.get(name)

    def get_skill_summary(self) -> Dict:
        """获取所有Skills的摘要"""
        return {
            "skills": {name: {
                "name": skill.name,
                "description": skill.description[:100] + "..." if len(skill.description) > 100 else skill.description,
                "operations": list(skill.operations.keys()),
                "keywords": skill.keywords
            } for name, skill in self.skills.items()},
            "agents": {name: {
                "name": agent.name,
                "description": agent.description,
                "use_cases": agent.use_cases[:3]
            } for name, agent in self.agents.items()}
        }

# ============ 语义分析器 ============

class SemanticAnalyzer:
    """基于LLM的语义分析器 - 完全使用LLM进行意图分析，不依赖关键词匹配"""

    def __init__(self, discovery: SkillDiscovery, model_caller):
        self.discovery = discovery
        self.model_caller = model_caller  # 传入模型调用函数
        self.decomposer = TaskDecomposer(model_caller, discovery)  # 将 discovery 传递给分解器

    def analyze(self, user_input: str, context: Optional[Dict] = None) -> Intent:
        """分析用户意图 - 先检测复杂请求，再进行语义分析"""
        logger.info(f"[Analyzer] 分析用户输入: {user_input[:50]}...")

        # 第1步：检测是否为复杂多步骤请求
        if self.decomposer.is_complex_request(user_input):
            logger.info(f"[Analyzer] 检测到复杂请求，进行任务分解...")
            model = context.get("model", "haiku") if context else "haiku"

            # 分解复杂请求
            metadata = self.decomposer.decompose(user_input, model=model)

            if metadata.is_complex:
                logger.info(f"[Analyzer] 分解完成: {metadata.task_count} 个子任务")

                # 创建包含任务链的Intent
                intent = Intent(
                    intent_type="complex_task_chain",
                    target="task_orchestrator",
                    is_complex=True,
                    subtasks=metadata.subtasks_list,
                    reasoning=f"复杂请求，分解为 {metadata.task_count} 个子任务"
                )
                return intent

        # 第2步：如果不是复杂请求，使用原有的LLM分析
        logger.info(f"[Analyzer] 简单请求，进行常规分析...")
        llm_intent = self._llm_analyze(user_input, context)
        return llm_intent
#非复杂，简单
    def _llm_analyze(self, user_input: str, context: Optional[Dict]) -> Intent:
        """使用LLM进行语义分析 - 完全基于LLM判断意图"""
        # 从context获取用户选择的模型
        model = context.get("model", "haiku") if context else "haiku"

        # 构建系统提示
        system_prompt = self._build_analysis_prompt()

        # 构建可用的 target 列表（包含 skills 和 agents）
        available_targets = ["WebSearch_dk","send_mail","claude"] #前两个是本地函数
        if self.discovery:
            if self.discovery.skills:
                available_targets.extend(self.discovery.skills.keys())
            if self.discovery.agents:
                available_targets.extend(self.discovery.agents.keys())
        target_str = "|".join(available_targets)

        user_prompt = f"""请分析以下用户请求，判断应该执行什么操作。

用户输入: {user_input}

请仔细分析用户意图，然后以 JSON 格式返回分析结果（只返回JSON，不要其他内容）。
** 重要：必须返回有效的JSON格式，内部字符串必须正确转义。**

{{
  "intent_type": "skill|direct|complex_task_chain|agent",
  "target": "{target_str}",
  "operation": "操作名称（见上方列表）",
  "params": {{"相关参数": "值"}},
  "reasoning": "简要说明判断理由"
}}
"""
# 如果是发送邮件，intent_type是direct, operation是send_mail;如果是网络搜索，intent_type是direct, operation是WebSearch_dk;""" #WebSearch(内置的，目前需要付费)

        try:
            logger.info(f"[Analyzer] 使用LLM分析意图，模型: {model}")
            # 使用两个参数调用 model_caller: system_prompt 和 user_prompt
            response = self.model_caller(system_prompt, user_prompt, model=model, temperature=0.1)
            logger.debug(f"[Analyzer] LLM响应: {response}")

            # 尝试解析 JSON
            result = self._extract_json(response)
            if result:
                logger.info(f"[Analyzer] LLM分析结果: type={result.get('intent_type')}, target={result.get('target')}, op={result.get('operation')}")
                # 只提取 Intent 类需要的字段，忽略多余字段
                intent_fields = {
                    'intent_type': result.get('intent_type'),
                    'target': result.get('target'),
                    'operation': result.get('operation'),
                    'params': result.get('params', {}),
                    'confidence': result.get('confidence', 1.0),
                    'reasoning': result.get('reasoning', ''),
                    'is_complex': result.get('is_complex', False),
                    'subtasks': result.get('subtasks', []),
                }
                # 移除值为 None 的字段，让 dataclass 使用默认值
                intent_fields = {k: v for k, v in intent_fields.items() if v is not None and v != {} and v != [] and v != False and v != ''}
                return Intent(**intent_fields)
        except Exception as e:
            logger.warning(f"[Analyzer] LLM 分析失败: {e}")

        # LLM 失败时，使用本地关键词规则兜底
        fallback_intent = self._local_fallback_rules(user_input)
        if fallback_intent:
            logger.info(f"[Analyzer] 使用本地规则兜底: type={fallback_intent.intent_type}, target={fallback_intent.target}, op={fallback_intent.operation}")
            return fallback_intent

        # 默认：直接让 Claude 处理
        return Intent(
            intent_type="direct",
            target="claude",
            reasoning="LLM 分析失败，使用默认处理"
        )

    def _local_fallback_rules(self, user_input: str) -> Optional[Intent]:
        """本地关键词规则兜底 - 当 LLM 意图分析失败时使用"""
        input_lower = user_input.lower()

#         # 0. 邮件发送（提取邮箱地址）- 最高优先级
#         email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
#         email_matches = re.findall(email_pattern, user_input)
#         if email_matches:
#             # 检测发送相关关键词
#             # send_keywords = ['发送', '发到', '发给', '邮件', '邮箱', 'email', '发']
#             # has_send = any(kw in input_lower for kw in send_keywords)

#             # 检测文档相关关键词
#             doc_keywords = ['文档', 'docx', 'word', 'ppt', 'pptx', 'md', 'markdown', '内容']
#             has_doc = any(kw in input_lower for kw in doc_keywords)

#             # 检测知识库相关关键词
#             kb_keywords = ['知识库', 'kb', '库中', '库里']
#             has_kb = any(kw in input_lower for kw in kb_keywords)

#             if has_send and (has_doc or has_kb):
#                 # 检测目标知识库
#                 kb = "KB"
#                 if '国学' in input_lower or 'kbgx' in input_lower:
#                     kb = "KBGX"
#                 elif '微信' in input_lower or 'kbw' in input_lower:
#                     kb = "KBW"

#                 # 检测文档格式
#                 doc_format = "docx"
#                 if 'ppt' in input_lower:
#                     doc_format = "pptx"
#                 elif 'md' in input_lower or 'markdown' in input_lower:
#                     doc_format = "md"
# #todo
#                 return Intent(
#                     intent_type="skill",
#                     target="knowledge-base",
#                     operation="email_send",
#                     params={"email": email_matches[0], "format": doc_format, "kb": kb},
#                     reasoning=f"本地规则: 检测到邮箱地址 {email_matches[0]} → email_send"
#                 )
#             elif has_send:
#                 # 仅发送内容到邮箱
#                 return Intent(
#                     intent_type="skill",
#                     target="knowledge-base",
#                     operation="email_send",
#                     params={"email": email_matches[0], "format": "text"},
#                     reasoning=f"本地规则: 发送内容到邮箱 {email_matches[0]}"
#                 )

        # 1. 保存/生成文档到知识库
        save_keywords = ['保存', '存储', '存到', '存入', '写入']
        doc_keywords = ['文档', 'docx', 'word', 'ppt', 'pptx', 'md', 'markdown']
        kb_keywords = ['知识库', 'kb', '库中', '库里']

        has_save = any(kw in input_lower for kw in save_keywords)
        has_kb = any(kw in input_lower for kw in kb_keywords)

        if has_save and has_kb:
            # 检测目标知识库
            kb = "KB"  # 默认
            if '国学' in input_lower or 'kbgx' in input_lower:
                kb = "KBGX"
            elif '微信' in input_lower or 'kbw' in input_lower:
                kb = "KBW"

            # 检测文档格式
            doc_format = "docx"  # 默认
            if 'ppt' in input_lower:
                doc_format = "pptx"
            elif 'md' in input_lower or 'markdown' in input_lower:
                doc_format = "md"

            return Intent(
                intent_type="skill",
                target="knowledge-base",
                operation="document_generation",
                params={"format": doc_format, "kb": kb},
                reasoning="本地规则: 保存到知识库 → document_generation"
            )

        # 2. 搜索知识库
        search_keywords = ['搜索', '查找', '查询', '检索', '找']
        if any(kw in input_lower for kw in search_keywords) and has_kb:
            kb = "KB"
            if '国学' in input_lower or 'kbgx' in input_lower:
                kb = "KBGX"
            elif '微信' in input_lower or 'kbw' in input_lower:
                kb = "KBW"
            return Intent(
                intent_type="skill",
                target="knowledge-base",
                operation="search",
                params={"kb": kb},
                reasoning="本地规则: 搜索知识库 → search"
            )

        # 3. 网络搜索
        web_keywords = ['天气', '新闻', '股票', '汇率', '网络搜索', 'web搜索']
        if any(kw in input_lower for kw in web_keywords):
            return Intent(
                intent_type="direct",
                target="WebSearch",
                reasoning="本地规则: 网络搜索关键词"
            )

        # 4. 列出知识库文件
        list_keywords = ['列出', '显示文件', '文件列表', '有哪些文件']
        if any(kw in input_lower for kw in list_keywords) and has_kb:
            kb = "KB"
            if '国学' in input_lower:
                kb = "KBGX"
            elif '微信' in input_lower:
                kb = "KBW"
            return Intent(
                intent_type="skill",
                target="knowledge-base",
                operation="list_docs",
                params={"kb": kb},
                reasoning="本地规则: 列出文件 → list_docs"
            )

        # 5. 下载上文内容
        context_download_keywords = ['上述', '上面', '刚才', '之前']
        download_keywords = ['下载', '保存为', '导出']
        if any(kw in input_lower for kw in context_download_keywords) and any(kw in input_lower for kw in download_keywords):
            doc_format = "md"
            if 'docx' in input_lower or 'word' in input_lower:
                doc_format = "docx"
            elif 'ppt' in input_lower:
                doc_format = "pptx"
            return Intent(
                intent_type="skill",
                target="knowledge-base",
                operation="context_download",
                params={"format": doc_format},
                reasoning="本地规则: 下载上文 → context_download"
            )

        # 没有匹配的本地规则
        return None

    def _build_analysis_prompt(self) -> str:
        """构建分析提示词 - 包含所有可用操作的详细说明"""
        prompt = """你是意图分析专家。请根据用户输入判断应该执行哪个操作。

# 可用操作列表(举例，非全列表)

## 1. 网络搜索 (WebSearch)
- **用途**: 搜索互联网获取实时信息
- **触发场景**: 用户询问天气、新闻、股票、实时价格、最新消息、当前状态等需要联网获取的信息
- **返回格式**: intent_type="direct", target="WebSearch_dk"
- **示例输入**:
  - "今天上海天气怎么样"
  - "最新的科技新闻"
  - "当前美元汇率"
  - "网络搜索xxx"

## 2. 知识库操作 (knowledge-base)

### 2.1 列出文件 (list_docs)
- **用途**: 显示知识库中的文件列表
- **触发场景**: 用户要查看知识库有哪些文件、文档目录、文件清单
- **返回格式**: intent_type="skill", target="knowledge-base", operation="list_docs"
- **参数**: params={"kb": "KB|KBGX|KBW"}
  - KB: 通用知识库
  - KBGX: 国学知识库（含道德经、论语、庄子等）
  - KBW: 微信知识库
- **示例输入**:
  - "列出知识库文件"
  - "显示国学文档"
  - "知识库有哪些文件"

### 2.2 显示全文 (show_full_doc)
- **用途**: 显示某个文档的完整内容
- **触发场景**: 用户要查看某个具体文档的全部内容
- **返回格式**: intent_type="skill", target="knowledge-base", operation="show_full_doc"
- **参数**: params={"doc_name": "文档名", "kb": "KB|KBGX|KBW"}
- **示例输入**:
  - "显示 华为IT供应商资质与需求 内容"
  - "查看金刚经主要智慧总结.md全文"

### 2.3 知识库搜索 (search)
- **用途**: 在知识库中搜索特定内容
- **触发场景**: 用户要在知识库中查找信息、搜索内容、查询知识
- **返回格式**: intent_type="skill", target="knowledge-base", operation="search"
- **参数**: params={"kb": "KB|KBGX|KBW", "query": "搜索词"}
- **示例输入**:
  - "搜索道德经中关于无为的内容"
  - "在知识库中查找关于修炼的信息"
  - "道德经第二章"

### 2.4 提取章节 (extract_chapter)
- **用途**: 提取文档的特定章节、页面或序号内容（支持EPUB和PDF）
- **触发场景**: 用户要查看某个文档的特定章、节、回、页或序号
- **返回格式**: intent_type="skill", target="knowledge-base", operation="extract_chapter"
- **参数**: params={"doc_name": "文档名", "chapter": "章节名/页码/序号", "kb": "KB|KBGX|KBW"}
- **支持的格式**:
  - 章节名：第一章、第二章、编者的话、上经、下经、第一部分等
  - 页码：第1页、第2页、第10页、1-10页
  - 序号：第一个、第二个、第1个、第2个 或 序号1、序号2
  - 中文数字：一、二、三、四、五、六、七、八、九、十
  - 阿拉伯数字：1、2、3...
- **支持的文件格式**: EPUB电子书、PDF文档
- **示例输入**:
  - "显示道德经第三章" (EPUB)
  - "显示A股进入牛市第一部分" (PDF)
  - "显示道德经第3页"
  - "显示道德经的第二个章节"
  - "查看道德经序号5的内容"
  - "显示道德经说什么的编者的话"

### 2.5 总结文档 (summarize_doc)
- **用途**: 对文档内容进行总结、提炼要点
- **触发场景**: 用户要求总结、归纳、提炼、概述某个文档
- **返回格式**: intent_type="skill", target="knowledge-base", operation="summarize_doc"
- **参数**: params={"doc_name": "文档名", "kb": "KB|KBGX|KBW"}
- **示例输入**:
  - "总结华为IT供应商资质与需求"
  - "归纳道德经的核心要义"

### 2.6 下载现有文件 (file_download)
- **用途**: 下载知识库中已存在的文件
- **触发场景**: 用户要下载一个具体的、带扩展名的文件
- **返回格式**: intent_type="skill", target="knowledge-base", operation="file_download"
- **参数**: params={"filename": "完整文件名.扩展名"}
- **示例输入**:
  - "下载 金刚经主要智慧总结.md"
  - "下载 华为培训.pptx"

### 2.7 下载上文内容 (context_download)
- **用途**: 将之前对话中的内容保存为文件下载
- **触发场景**: 用户引用"上述"、"上面"、"刚才"的内容要求保存/下载
- **返回格式**: intent_type="skill", target="knowledge-base", operation="context_download"
- **参数**: params={"format": "md|docx|pptx"}
- **示例输入**:
  - "把上面的内容保存为md下载"
  - "将刚才的回答转为word下载"

### 2.8 生成文档 (document_generation)
- **用途**: 基于内容生成新的文档文件（Word、PPT、Markdown等）
- **触发场景**: 用户要求将内容转换为特定格式的文档，或生成/创建文档
- **返回格式**: intent_type="skill", target="knowledge-base", operation="document_generation"
- **参数**: params={"format": "docx|pptx|md|pdf"}
- **示例输入**:
  - "把这些内容转为PPT"
  - "生成Word文档"
  - "创建一个关于xxx的演示文稿"
  
## 3. 发送邮件 (send_mail)
- **用途**: 将输入或上下文中内容、文档链接发送到邮箱
- **触发场景**: 用户要求发邮件
- **返回格式**: intent_type="direct", target="send_mail"
- **示例输入**:
  - "将下面内容发送到xxx@example.com:xxxx"
  - "将上面内容（或生成的文件）发到xxx@example.com:xxxx"

## 4. Bash操作 (Bash)
- **用途**: github代码库访问（通过Bash）、查看当前文件夹文件等
- **触发场景**: 要求访问代码库或本地文件夹列表等
- **返回格式**: intent_type="agent", target="Bash"
- **示例输入**:
  - "你好"
  - "解释一下量子力学"
   
## 4. 普通对话 (claude)
- **用途**: 普通的问答、聊天、解释、帮助等
- **触发场景**: 不属于上述任何明确操作的请求
- **返回格式**: intent_type="direct", target="claude",operation="Bash|general-purpose|statusline-setup|Explore|Plan"
- **示例输入**:
  - "你好"
  - "解释一下量子力学"

# 知识库识别规则
- **KBGX (国学知识库)**: 涉及道德经、论语、庄子、金刚经、易经、诗经、孟子、大学、中庸、六祖坛经等中国古典著作
- **KBW (微信知识库)**: 涉及微信、公众号、订阅号、服务号相关内容
- **KB (通用知识库)**: 其他文档

# 判断优先级
1. 明确的网络搜索请求（天气、新闻、实时信息）→ WebSearch
2. 明确的文件下载请求（下载 + 完整文件名.扩展名）→ file_download
3. 引用上文的下载/保存请求 → context_download
4. 文档格式转换请求 → document_generation
5. 知识库文件列表请求 → list_docs
6. 显示全文请求 → show_full_doc
7. 章节提取请求 → extract_chapter
8. 总结请求 → summarize_doc
9. 知识库搜索请求 → search
10. 发送邮件 → send_mail 
10. 其他 → claude

请只返回JSON，不要添加任何解释文字。""" # 加上bash，webfetch,websearch
        return prompt
    def _extract_json(self, text: str) -> Optional[Dict]: #deepseek的方案好，另外一个也要换
        """从文本中提取 JSON，优先获取最外层的完整对象"""
        if not text:
            return None

        # 0. 先移除 <think>...</think> 标签（DeepSeek/MiniMax等模型可能返回带思考过程的响应）
        text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL).strip()

        if not text:
            return None

        def try_parse_json(json_str: str) -> Optional[Dict]:
            """尝试解析JSON字符串"""
            try:
                return json.loads(json_str)
            except json.JSONDecodeError as e:
                logger.debug(f"[Analyzer] JSON 解析失败: {e}")
                return None

        # 1. 尝试直接解析整个文本
        result = try_parse_json(text.strip())
        if result is not None:
            return result

        # 2. 尝试从 markdown 代码块中提取 JSON (```json ... ```)
        match = re.search(r'```(?:json)?\s*\n(.*?)\n```', text, re.DOTALL)
        if match:
            json_str = match.group(1).strip()
            result = try_parse_json(json_str)
            if result is not None:
                logger.debug(f"[Analyzer] 从 markdown 代码块成功提取 JSON")
                return result
            
            logger.debug(f"[Analyzer] 从代码块解析JSON失败，尝试修复JSON字符串")
            # 修复JSON字符串中的常见问题
            try:
                # 修复未转义的双引号：将字符串内部的双引号转义，但不转义键名和属性名周围的双引号
                # 通过正则表达式匹配字符串值部分
                def escape_inner_quotes(match):
                    # 匹配JSON字符串值（不包括键名）
                    content = match.group(1)
                    # 转义内部的双引号
                    content = content.replace('"', '\\"')
                    return f'"{content}"'
                
                # 匹配JSON字符串值：从冒号或逗号后的空格开始到下一个逗号、大括号或括号
                # 这个正则表达式更精确地匹配JSON字符串值
                json_str_fixed = re.sub(r'(?<=\:\s*)"(.*?)"(?=\s*[,}\]])', escape_inner_quotes, json_str)
                
                # 尝试解析修复后的JSON
                result = try_parse_json(json_str_fixed)
                if result is not None:
                    logger.debug(f"[Analyzer] 代码块JSON修复成功")
                    return result
            except Exception as e:
                logger.debug(f"[Analyzer] 代码块JSON修复失败: {e}")

        # 3. 最可靠的方法：找第一个 { 和最后一个 }（这是最外层的对象）
        first_brace = text.find('{')
        last_brace = text.rfind('}')
        if first_brace != -1 and last_brace != -1 and first_brace < last_brace:
            json_str = text[first_brace:last_brace+1]
            result = try_parse_json(json_str)
            if result is not None:
                logger.debug(f"[Analyzer] 从第一个和最后一个大括号成功提取 JSON")
                return result
            
            logger.debug(f"[Analyzer] JSON 解析失败，尝试智能修复...")
            try:
                # 智能修复JSON字符串
                lines = json_str.split('\n')
                fixed_lines = []
                
                for i, line in enumerate(lines):
                    line = line.rstrip()
                    if not line:
                        continue
                    
                    # 检查是否缺少逗号分隔符（当前行以"结尾，下一行以"开头但不是"}]）
                    if i < len(lines) - 1:
                        next_line = lines[i + 1].strip()
                        if line.endswith('"') and next_line.startswith('"'):
                            # 在当前行末尾添加逗号
                            line = line + ','
                    
                    fixed_lines.append(line)
                
                # 重新组合
                json_str_fixed = '\n'.join(fixed_lines)
                
                # 最后尝试：转义字符串内部的双引号
                def escape_string_values(match):
                    content = match.group(1)
                    # 跳过已经转义的双引号
                    if '\\"' in content:
                        return match.group(0)
                    # 转义未转义的双引号
                    content = content.replace('"', '\\"')
                    return f'"{content}"'
                
                # 匹配JSON字符串值（更精确的匹配）
                json_str_fixed = re.sub(r'(?<=\:\s*)"([^"\\]*(?:\\.[^"\\]*)*)"', escape_string_values, json_str_fixed)
                
                result = try_parse_json(json_str_fixed)
                if result is not None:
                    logger.debug(f"[Analyzer] JSON 智能修复成功")
                    return result
            except Exception as e:
                logger.debug(f"[Analyzer] JSON 智能修复失败: {e}")

                logger.warning(f"[Analyzer] 无法从文本中提取有效的 JSON")
        # 最后尝试：返回一个包含原始文本的简单JSON结构
        #不需要
        # try:
        #     return {
        #         "intent_type": "direct",
        #         "target": "send_mail",
        #         "operation": "send_mail",
        #         "params": {
        #             "email": "hanwsf@163.com",
        #             "content": text[text.find('{'):text.rfind('}')+1] if text.find('{') != -1 and text.rfind('}') != -1 else text
        #         }
        #     }
        # except:
        #     return None
    # def _extract_json(self, text: str) -> Optional[Dict]:
    #     """从文本中提取 JSON，优先获取最外层的完整对象"""
    #     if not text:
    #         return None

    #     # 1. 尝试直接解析整个文本
    #     try:
    #         return json.loads(text.strip())
    #     except:
    #         pass

    #     # 2. 尝试从 markdown 代码块中提取 JSON (```json ... ```)
    #     match = re.search(r'```(?:json)?\s*\n(.*?)\n```', text, re.DOTALL)
    #     if match:
    #         try:
    #             json_str = match.group(1).strip()
    #             result = json.loads(json_str)
    #             logger.debug(f"[Analyzer] 从 markdown 代码块成功提取 JSON")
    #             return result
    #         except Exception as e:
    #             logger.debug(f"[Analyzer] 从代码块解析JSON失败: {e}")
    #             # 代码块内容可能包含真实的换行符，尝试修复
    #             try:
    #                 json_str_fixed = json_str.replace('\n', '\\n').replace('\r', '\\r')
    #                 result = json.loads(json_str_fixed)
    #                 logger.debug(f"[Analyzer] 代码块JSON修复成功")
    #                 return result
    #             except Exception as e2:
    #                 logger.debug(f"[Analyzer] 代码块JSON修复也失败: {e2}")

    #     # 3. 最可靠的方法：找第一个 { 和最后一个 }（这是最外层的对象）
    #     first_brace = text.find('{')
    #     last_brace = text.rfind('}')
    #     if first_brace != -1 and last_brace != -1 and first_brace < last_brace:
    #         try:
    #             json_str = text[first_brace:last_brace+1]
    #             result = json.loads(json_str)
    #             logger.debug(f"[Analyzer] 从第一个和最后一个大括号成功提取 JSON")
    #             return result
    #         except json.JSONDecodeError as e:
    #             logger.debug(f"[Analyzer] JSON 解析失败: {e}, 尝试修复...")
    #             # 尝试通过替换可能的问题来修复 JSON
    #             try:
    #                 # 替换未转义的换行符和回车符
    #                 json_str_fixed = json_str.replace('\n', '\\n').replace('\r', '\\r')
    #                 result = json.loads(json_str_fixed)
    #                 logger.debug(f"[Analyzer] JSON 修复成功")
    #                 return result
    #             except Exception as e2:
    #                 logger.debug(f"[Analyzer] JSON 修复失败: {e2}")

    #     logger.warning(f"[Analyzer] 无法从文本中提取有效的 JSON")
    #     return None

# ============ 调度器 ============
#   覆盖的意图类型：                                                                                                                                 
#   1. WebSearch - 网络搜索                                                                                                                         
#   2. list_docs - 列出文件                                                                                                                          
#   3. show_full_doc - 显示全文                                                                                                                      
#   4. search - 知识库搜索                                                                                                                           
#   5. extract_chapter - 提取章节                                                                                                                    
#   6. summarize_doc - 总结文档                                                                                                                      
#   7. file_download - 下载现有文件                                                                                                                  
#   8. context_download - 下载上文内容                                                                                                               
#   9. document_generation - 生成文档                                                                                                                
#   10. claude - 普通对话        

class Dispatcher:
    """调度器 - 根据意图执行相应操作"""

    def __init__(self, discovery: SkillDiscovery, model_caller, kb_paths: Optional[Dict] = None):
        self.discovery = discovery
        self.model_caller = model_caller
        self.kb_paths = kb_paths or {}

    def dispatch(self, intent: Intent, user_input: str, context: Optional[Dict] = None) -> str:
        """执行调度"""
        logger.info(f"[Dispatcher] 调度 {intent.intent_type}: {intent.target}")

        # 处理复杂任务链
        if intent.intent_type == "complex_task_chain":
            return self._dispatch_complex_task_chain(intent, user_input, context)
        elif intent.intent_type == "skill":
            return self._dispatch_skill(intent, user_input, context)
        elif intent.intent_type == "agent":
            return self._dispatch_agent(intent, user_input, context)
        elif intent.intent_type == "direct" and intent.target == "WebSearch_dk":
            # 使用 dk-search 服务进行网络搜索
            return self._dispatch_web_search_dk(user_input, context)
        elif intent.intent_type == "direct" and intent.target == "WebFetch":
            # 获取网页内容（如 GitHub URL）
            return self._dispatch_web_fetch(user_input, context)
        elif intent.intent_type == "direct" and intent.target == "WebSearch":
            # 使用 Claude 内置 WebSearch 工具（保留原有功能）
            return self._dispatch_web_search(user_input, context)
        elif intent.intent_type == "direct" and intent.target == "send_mail":
            # 特殊处理网络搜索 
            return self._dispatch_send_mail(user_input, context)
        else: #claude
            return self._dispatch_direct(user_input, context)
    # def _extract_json(self, text: str) -> Optional[Dict]:
    #     """从文本中提取 JSON，优先获取最外层的完整对象"""
    #     if not text:
    #         return None

    #     # 1. 尝试直接解析整个文本
    #     try:
    #         return json.loads(text.strip())
    #     except:
    #         pass

    #     # 2. 尝试从 markdown 代码块中提取 JSON (```json ... ```)
    #     match = re.search(r'```(?:json)?\s*\n(.*?)\n```', text, re.DOTALL)
    #     if match:
    #         try:
    #             json_str = match.group(1).strip()
    #             result = json.loads(json_str)
    #             logger.debug(f"[Dispatcher] 从 markdown 代码块成功提取 JSON")
    #             return result
    #         except Exception as e:
    #             logger.debug(f"[Dispatcher] 从代码块解析JSON失败: {e}")
    #             # 代码块内容可能包含真实的换行符，尝试修复
    #             try:
    #                 json_str_fixed = json_str.replace('\n', '\\n').replace('\r', '\\r')
    #                 result = json.loads(json_str_fixed)
    #                 logger.debug(f"[Dispatcher] 代码块JSON修复成功")
    #                 return result
    #             except Exception as e2:
    #                 logger.debug(f"[Dispatcher] 代码块JSON修复也失败: {e2}")

    #     # 3. 最可靠的方法：找第一个 { 和最后一个 }（这是最外层的对象）
    #     first_brace = text.find('{')
    #     last_brace = text.rfind('}')
    #     if first_brace != -1 and last_brace != -1 and first_brace < last_brace:
    #         try:
    #             json_str = text[first_brace:last_brace+1]
    #             result = json.loads(json_str)
    #             logger.debug(f"[Dispatcher] 从第一个和最后一个大括号成功提取 JSON")
    #             return result
    #         except json.JSONDecodeError as e:
    #             logger.debug(f"[Dispatcher] JSON 解析失败: {e}, 尝试修复...")
    #             # 尝试通过替换可能的问题来修复 JSON
    #             try:
    #                 # 替换未转义的换行符和回车符
    #                 json_str_fixed = json_str.replace('\n', '\\n').replace('\r', '\\r')
    #                 result = json.loads(json_str_fixed)
    #                 logger.debug(f"[Dispatcher] JSON 修复成功")
    #                 return result
    #             except Exception as e2:
    #                 logger.debug(f"[Dispatcher] JSON 修复失败: {e2}")

    #     logger.warning(f"[Dispatcher] 无法从文本中提取有效的 JSON")
    #     return None
    def _extract_json(self, text: str) -> Optional[Dict]: #deepseek的方案好，另外一个也要换
        """从文本中提取 JSON，优先获取最外层的完整对象"""
        if not text:
            return None

        # 0. 先移除 <think>...</think> 标签（DeepSeek/MiniMax等模型可能返回带思考过程的响应）
        text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL).strip()

        if not text:
            return None

        def try_parse_json(json_str: str) -> Optional[Dict]:
            """尝试解析JSON字符串"""
            try:
                return json.loads(json_str)
            except json.JSONDecodeError as e:
                logger.debug(f"[Analyzer] JSON 解析失败: {e}")
                return None

        # 1. 尝试直接解析整个文本
        result = try_parse_json(text.strip())
        if result is not None:
            return result

        # 2. 尝试从 markdown 代码块中提取 JSON (```json ... ```)
        match = re.search(r'```(?:json)?\s*\n(.*?)\n```', text, re.DOTALL)
        if match:
            json_str = match.group(1).strip()
            result = try_parse_json(json_str)
            if result is not None:
                logger.debug(f"[Analyzer] 从 markdown 代码块成功提取 JSON")
                return result
            
            logger.debug(f"[Analyzer] 从代码块解析JSON失败，尝试修复JSON字符串")
            # 修复JSON字符串中的常见问题
            try:
                # 修复未转义的双引号：将字符串内部的双引号转义，但不转义键名和属性名周围的双引号
                # 通过正则表达式匹配字符串值部分
                def escape_inner_quotes(match):
                    # 匹配JSON字符串值（不包括键名）
                    content = match.group(1)
                    # 转义内部的双引号
                    content = content.replace('"', '\\"')
                    return f'"{content}"'
                
                # 匹配JSON字符串值：从冒号或逗号后的空格开始到下一个逗号、大括号或括号
                # 这个正则表达式更精确地匹配JSON字符串值
                json_str_fixed = re.sub(r'(?<=\:\s*)"(.*?)"(?=\s*[,}\]])', escape_inner_quotes, json_str)
                
                # 尝试解析修复后的JSON
                result = try_parse_json(json_str_fixed)
                if result is not None:
                    logger.debug(f"[Analyzer] 代码块JSON修复成功")
                    return result
            except Exception as e:
                logger.debug(f"[Analyzer] 代码块JSON修复失败: {e}")

        # 3. 最可靠的方法：找第一个 { 和最后一个 }（这是最外层的对象）
        first_brace = text.find('{')
        last_brace = text.rfind('}')
        if first_brace != -1 and last_brace != -1 and first_brace < last_brace:
            json_str = text[first_brace:last_brace+1]
            result = try_parse_json(json_str)
            if result is not None:
                logger.debug(f"[Analyzer] 从第一个和最后一个大括号成功提取 JSON")
                return result
            
            logger.debug(f"[Analyzer] JSON 解析失败，尝试智能修复...")
            try:
                # 智能修复JSON字符串
                lines = json_str.split('\n')
                fixed_lines = []
                
                for i, line in enumerate(lines):
                    line = line.rstrip()
                    if not line:
                        continue
                    
                    # 检查是否缺少逗号分隔符（当前行以"结尾，下一行以"开头但不是"}]）
                    if i < len(lines) - 1:
                        next_line = lines[i + 1].strip()
                        if line.endswith('"') and next_line.startswith('"'):
                            # 在当前行末尾添加逗号
                            line = line + ','
                    
                    fixed_lines.append(line)
                
                # 重新组合
                json_str_fixed = '\n'.join(fixed_lines)
                
                # 最后尝试：转义字符串内部的双引号
                def escape_string_values(match):
                    content = match.group(1)
                    # 跳过已经转义的双引号
                    if '\\"' in content:
                        return match.group(0)
                    # 转义未转义的双引号
                    content = content.replace('"', '\\"')
                    return f'"{content}"'
                
                # 匹配JSON字符串值（更精确的匹配）
                json_str_fixed = re.sub(r'(?<=\:\s*)"([^"\\]*(?:\\.[^"\\]*)*)"', escape_string_values, json_str_fixed)
                
                result = try_parse_json(json_str_fixed)
                if result is not None:
                    logger.debug(f"[Analyzer] JSON 智能修复成功")
                    return result
            except Exception as e:
                logger.debug(f"[Analyzer] JSON 智能修复失败: {e}")

        logger.warning(f"[Analyzer] 无法从文本中提取有效的 JSON")
        # 最后尝试：返回一个包含原始文本的简单JSON结构
        try:
            return {
                "intent_type": "direct",
                "target": "send_mail",
                "operation": "send_mail",
                "params": {
                    "email": "hanwsf@163.com",
                    "content": text[text.find('{'):text.rfind('}')+1] if text.find('{') != -1 and text.rfind('}') != -1 else text
                }
            }
        except:
            return None
    def _build_prior_results_context(self, task: SubTask, task_results: Dict) -> str:
        """构建前置步骤结果的上下文信息"""
        if not task.depends_on or not task_results:
            return ""

        parts = []
        for dep_id in task.depends_on:
            if dep_id in task_results:
                result_text = str(task_results[dep_id])
                # 限制每个步骤结果的长度，避免 prompt 过长
                if len(result_text) > 16000:
                    result_text = result_text[:16000] + "\n...(内容过长已截断)"
                parts.append(f"【{dep_id} 的输出结果】:\n{result_text}")

        if not parts:
            return ""

        return "\n\n=== 前置步骤已完成的结果（请基于这些结果继续工作） ===\n" + "\n\n".join(parts) + "\n=== 前置步骤结果结束 ===\n"

    def _dispatch_complex_task_chain(self, intent: Intent, user_input: str, context: Optional[Dict]) -> str:
        """调度复杂任务链 - 执行多个子任务"""
        #todo 完善
        logger.info(f"[Dispatcher] 执行复杂任务链: {len(intent.subtasks)} 个子任务,intent:{str(intent)}")

        if not intent.subtasks:
            return "❌ 复杂任务链为空"

        # 定义子任务执行函数
        def execute_subtask(task: SubTask, task_results: Dict) -> Any:
            """执行单个子任务"""
            logger.info(f"[SubTask] 执行 {task.task_id}: {task.name} (目标: {task.target})")

            try:
                # 根据任务类型调度
                if task.intent_type == "skill": #docx
                    return self._execute_subtask_skill(task, user_input, context, task_results)
                elif task.intent_type == "agent":
                    return self._execute_subtask_agent(task, user_input, context, task_results)
                elif task.intent_type == "direct": #mail
                    return self._execute_subtask_direct(task, user_input, context, task_results)
                else:
                    return f"❌ 未知任务类型: {task.intent_type}"
            except Exception as e:
                logger.error(f"[SubTask] {task.task_id} 执行异常: {e}")
                raise

        # 执行任务链（同步版本）
        task_results = {}
        execution_order = self._topological_sort(intent.subtasks)

        for task in execution_order:
            # 检查依赖是否都已完成
            if all(dep in task_results for dep in task.depends_on):
                try:
                    logger.info(f"[SubTask] {task.task_id} 开始执行,{str(task)}")
                    task.start_time = datetime.now()
                    task.status = "running"

                    # 执行任务
                    result = execute_subtask(task, task_results)

                    task.result = result
                    task.status = "completed"
                    task.end_time = datetime.now()
                    task_results[task.task_id] = result

                    elapsed = (task.end_time - task.start_time).total_seconds()
                    logger.info(f"[SubTask] {task.task_id} 完成，耗时 {elapsed:.1f}秒,\n{result}")
                except Exception as e:
                    task.status = "failed"
                    task.error = str(e)
                    task.end_time = datetime.now()
                    logger.error(f"[SubTask] {task.task_id} 失败: {e}")
            else:
                # 依赖失败，标记为被阻止
                task.status = "blocked"
                task.error = "依赖的前置任务失败"
                logger.warning(f"[SubTask] {task.task_id} 被阻止（依赖失败）")

        # 构建执行结果总结
        completed_count = sum(1 for t in intent.subtasks if t.status == "completed")

        summary = self._build_task_chain_summary(intent.subtasks, task_results)

        logger.info(f"[Dispatcher] 任务链执行完成: {completed_count}/{len(intent.subtasks)} 成功")

        return summary

    def _topological_sort(self, tasks: List[SubTask]) -> List[SubTask]:
        """对任务进行拓扑排序"""
        sorted_tasks = []
        visited = set()
        in_progress = set()

        def visit(task_id: str):
            if task_id in visited:
                return
            if task_id in in_progress:
                raise ValueError(f"检测到循环依赖: {task_id}")

            in_progress.add(task_id)

            # 找到该任务
            task = next((t for t in tasks if t.task_id == task_id), None)
            if not task:
                return

            # 先访问依赖
            for dep in task.depends_on:
                visit(dep)

            in_progress.remove(task_id)
            visited.add(task_id)
            sorted_tasks.append(task)

        for task in tasks:
            visit(task.task_id)

        return sorted_tasks

    def _execute_subtask_skill(self, task: SubTask, user_input: str, context: Optional[Dict],task_results: Dict = None) -> str:
        """执行 Skill 类型的子任务"""
        skill = self.discovery.get_skill_by_name(task.target)
        logger.info(f"==sub task skill:{str(task)}")
        if not skill:
            # Fallback：使用 Claude 直接处理（+5分钟超时）
            logger.warning(f"[SubTask] Skill '{task.target}' 未找到，使用 Claude 处理（+5分钟超时）")
            task.timeout_seconds = 600  # 增加到 10 分钟
            return self._execute_subtask_direct(task, user_input, context)

        # 如果是知识库 skill，直接调用
        if skill.name == "knowledge-base":
            intent = Intent(
                intent_type="skill",
                target="knowledge-base",
                operation=task.operation,
                params=task.params
            )
            result = self._dispatch_knowledge_base(intent, user_input, context)
            logger.info(f"==========sub task skill {task.task_id}知识库结果:{result}")
            return result

        # # === 特殊处理 docx skill - 实际生成 Word 文档 ===
        # if skill.name == "docx":
        #     return self._handle_docx_skill(task, user_input, context, task_results)

        # 获取前置步骤的结果上下文
        prior_context = self._build_prior_results_context(task, task_results) if task_results else ""

        # 其他 skill，通过 LLM 调用
        prompt = f"""请使用 {skill.name} Skill 来处理以下用户请求。

用户请求: {user_input}
子任务名: {task.name}

Skill 描述: {skill.description}
{prior_context}

重要：如果前置步骤已经生成了内容（如文章、文档），请基于那些内容继续工作，而不是重新生成。"""

        model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"
        result = self.model_caller(prompt, model=model, temperature=0.3)
        logger.info(f"==========sub task skill {task.task_id}结果:{result}")
        return result

    def _execute_subtask_agent(self, task: SubTask, user_input: str, context: Optional[Dict], task_results: Dict = None) -> str:
        """执行 Agent 类型的子任务"""
        agent = self.discovery.get_agent_by_name(task.target)
        if not agent:
            # Fallback：使用 Claude 直接处理（+5分钟超时）
            logger.warning(f"[SubTask] Agent '{task.target}' 未找到，使用 Claude 处理（+5分钟超时）")
            task.timeout_seconds = 600  # 增加到 10 分钟
            return self._execute_subtask_direct(task, user_input, context)
        # 获取前置步骤的结果上下文
        prior_context = self._build_prior_results_context(task, task_results) if task_results else ""
        prompt = f"""请使用 {agent.name} Agent 来处理以下用户请求。

Agent 描述: {agent.description}
适用场景: {'; '.join(agent.use_cases)}
子任务: {task.name}

用户请求: {user_input}
{prior_context}

重要：如果前置步骤已经生成了内容（如文章、文档），请基于那些内容继续工作，而不是重新生成。"""

        model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"
        result = self.model_caller(prompt, model=model, temperature=0.3)
        logger.info(f"==========sub task skill {task.task_id}结果:{result}")
        return result
 
    def _validate_single_step_output(self, task: SubTask, output: str) -> str:                                                                                                                
        """验证输出是否只包含单步结果，防止假完成"""                                                                                                                                          
        import re                                                                                                                                                                             
                                                                                                                                                                                                
        # 检测多步骤标记                                                                                                                                                                      
        step_patterns = [                                                                                                                                                                     
            r'##\s*步骤\s*\d+',                                                                                                                                                               
            r'步骤\s*\d+[:：]',                                                                                                                                                               
            r'Step\s*\d+[:：]',                                                                                                                                                               
            r'第\s*\d+\s*步[:：]',                                                                                                                                                            
        ]                                                                                                                                                                                     
                                                                                                                                                                                                
        step_count = 0                                                                                                                                                                        
        for pattern in step_patterns:                                                                                                                                                         
            matches = re.findall(pattern, output, re.IGNORECASE)                                                                                                                              
            step_count += len(matches)                                                                                                                                                        
                                                                                                                                                                                                
        if step_count > 1:                                                                                                                                                                    
            logger.warning(f"[SubTask] {task.task_id} 输出包含 {step_count} 个步骤标记，可能是假完成")                                                                                        
                                                                                                                                                                                                
            # 尝试只提取第一个步骤的内容                                                                                                                                                      
            first_step_end = None                                                                                                                                                             
            for pattern in step_patterns:                                                                                                                                                     
                match = re.search(pattern, output[100:], re.IGNORECASE)  # 跳过开头                                                                                                           
                if match:                                                                                                                                                                     
                    if first_step_end is None or match.start() < first_step_end:                                                                                                              
                        first_step_end = match.start() + 100                                                                                                                                  
                                                                                                                                                                                                
            if first_step_end and first_step_end < len(output) - 200:                                                                                                                         
                logger.info(f"[SubTask] 截取第一个步骤的内容，原始长度 {len(output)}，截取后 {first_step_end}")                                                                               
                output = output[:first_step_end].strip()                                                                                                                                      
                output += "\n\n⚠️ [系统提示：检测到多步骤输出，已截取当前步骤内容]"                                                                                                           
                                                                                                                                                                                                
        # 检测是否包含后续任务的代码（如邮件发送）                                                                                                                                            
        dangerous_patterns = [                                                                                                                                                                
            (r'def send_email', '邮件发送代码'),                                                                                                                                              
            (r'smtplib\.SMTP', 'SMTP 邮件代码'),                                                                                                                                              
            (r'MIMEMultipart', '邮件附件代码'),                                                                                                                                               
            (r'def create_word_document.*def send', 'Word+邮件组合代码'),                                                                                                                     
        ]                                                                                                                                                                                     
                                                                                                                                                                                                
        for pattern, desc in dangerous_patterns:                                                                                                                                              
            if re.search(pattern, output, re.DOTALL):                                                                                                                                         
                logger.warning(f"[SubTask] {task.task_id} 输出包含不应该有的 {desc}")                                                                                                         
                # 移除这些代码                                                                                                                                                                
                output = re.sub(rf'{pattern}.*?(?=\ndef |\nclass |\n#\s*\n|$)',                                                                                                               
                                f'\n# [{desc} 已移除 - 将在后续步骤执行]\n',                                                                                                                   
                                output, flags=re.DOTALL)                                                                                                                                       
                                                                                                                                                                                                
        return output                                         
        
    def _execute_subtask_direct(self, task: SubTask, user_input: str,                                                                                                                         
                                context: Optional[Dict], task_results: Dict) -> str:                                                                                                          
        """执行 direct 类型的子任务 - 强化单步约束"""                                                                                                                                         
                                                                                                                                                                                                
        # 获取前置任务结果                                                                                                                                                                    
        previous_results = ""                                                                                                                                                                 
        for dep_id in (task.depends_on or []):                                                                                                                                                
            if dep_id in task_results:                                                                                                                                                        
                result_preview = str(task_results[dep_id])[:2000]                                                                                                                             
                previous_results += f"\n【{dep_id}的结果】:\n{result_preview}\n"                                                                                                              
                                                                                                                                                                                                
        # 强化约束的提示词                                                                                                                                                                    
        prompt = f"""你正在执行多步任务链中的 **第 {task.task_id} 步**（共多步）。                                                                                                            
                                                                                                                                                                                                
    ═══════════════════════════════════════════════════════════                                                                                                                               
    【当前任务】: {task.name}                                                                                                                                                                 
    【任务目标】: {task.operation if hasattr(task, 'operation') else task.name}                                                                                                               
    ═══════════════════════════════════════════════════════════                                                                                                                               
                                                                                                                                                                                                
    【严格约束 - 违反将导致任务失败】:                                                                                                                                                        
    1. ⚠️ 只执行当前这一步，不要执行或描述后续步骤                                                                                                                                            
    2. ⚠️ 不要生成其他步骤的代码（如邮件发送、文件转换等）                                                                                                                                    
    3. ⚠️ 如果任务是"获取内容"，只返回获取到的内容                                                                                                                                            
    4. ⚠️ 不要假装已完成后续任务                                                                                                                                                              
    5. ⚠️ 如果需要执行代码才能完成，生成可执行的代码块                                                                                                                                        
                                                                                                                                                                                                
    【用户原始请求】:                                                                                                                                                                         
    {user_input}                                                                                                                                                                              
                                                                                                                                                                                                
    【前置步骤结果】:                                                                                                                                                                         
    {previous_results if previous_results else "（这是第一步，无前置结果）"}                                                                                                                  
                                                                                                                                                                                                
    【输出格式】:                                                                                                                                                                             
    直接返回当前步骤的实际执行结果，不要包含步骤标记如"步骤1"、"步骤2"等。                                                                                                                    
    """                                                                                                                                                                                       
                                                                                                                                                                                                
        model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"                                                                         
        result = self.model_caller(prompt, model=model, temperature=0.3)                                                                                                                      
        logger.info(f"====================sub task {task.task_id} 输出结果：{result}")                                                                                                                                                                                        
        # 验证输出是否包含多步骤（可能是假完成）                                                                                                                                              
        result = self._validate_single_step_output(task, result)                                                                                                                              
        logger.info(f"====================sub task {task.task_id} 验证后输出结果：{result}")                                                                                                                                                                                         
        return result                                          
#     def _execute_subtask_direct(self, task: SubTask, user_input: str, context: Optional[Dict], task_results: Dict = None) -> str:
#         """执行 Direct 类型的子任务（直接调用 Claude）"""
#         logger.info(f"执行 Direct 类型的子任务（直接调用 Claude）,现在任务是：{str(task)}")
#         # === 检测并执行邮件发送任务 ===
#         # email_keywords = ['发送邮件', '邮件发送', 'send email', '发到邮箱', '发送到邮箱', '发邮件']
#         # task_name_lower = task.name.lower()
#         # is_email_task = any(kw in task_name_lower for kw in email_keywords)

#         # if is_email_task:
#         #     return self._handle_email_subtask(task, user_input, context, task_results)

#         # 获取前置步骤的结果上下文
#         prior_context = self._build_prior_results_context(task, task_results) if task_results else ""

#         prompt = f"""请处理以下用户请求的子任务。

# 任务名: {task.name}
# 用户请求: {user_input}
# {prior_context}

# 重要：
# 1. 如果前置步骤已经生成了内容（如完整的文章、文档），请直接使用那些内容，不要重新生成摘要版本。
# 2. 如果任务是发送邮件，请使用前置步骤生成的完整内容作为邮件正文或附件内容。
# 3. 文件名应该使用有意义的名称，不要使用通用名称。

# 请直接执行这个任务，不需要额外的确认。"""

#         model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"
#         return self.model_caller(prompt, model=model, temperature=0.3)
    #sub-task用
#     def _handle_email_subtask(self, task: SubTask, user_input: str,
#                               context: Optional[Dict], task_results: Dict) -> str:
#         """处理邮件发送子任务"""
#         logger.info(f"[SubTask] 检测到邮件发送任务: {task.name}")
# #估计这里需要处理上下文中的邮件内容
#         # 1. 从 task.params 或 user_input 提取邮箱地址
#         email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
#         to_email = task.params.get('email', '') if task.params else ''
#         if not to_email:
#             matches = re.findall(email_pattern, user_input)
#             to_email = matches[0] if matches else ''

#         if not to_email:
#             return "❌ 未找到收件人邮箱地址"

#         # 2. 从前置任务结果获取内容/附件
#         content = ""
#         attachment_path = None

#         for dep_id in (task.depends_on or []):
#             if dep_id in task_results:
#                 dep_result = task_results[dep_id]
#                 # 检查是否有文件路径
#                 if isinstance(dep_result, str):
#                     # 查找文件路径 - 优先查找本地生成的文件
#                     # 先检查是否有生成的 .docx 文件
#                     docx_match = re.search(r'([^\s]+\.docx)', dep_result)
#                     if docx_match:
#                         potential_path = docx_match.group(1)
#                         if os.path.exists(potential_path):
#                             attachment_path = potential_path
#                             logger.info(f"[Email] 检测到生成的 Word 文档: {potential_path}")

#                     # 其次查找其他格式的文件
#                     if not attachment_path:
#                         path_match = re.search(r'(/[^\s]+\.(docx|pptx|pdf|md|txt))', dep_result)
#                         if path_match and os.path.exists(path_match.group(1)):
#                             attachment_path = path_match.group(1)

#                     # 如果还没有内容，将结果作为邮件内容
#                     if not content:
#                         content = dep_result

#         # 3. 构建邮件主题
#         subject = task.params.get('subject', '') if task.params else ''
#         if not subject:
#             # 尝试从内容中提取标题
#             title_match = re.search(r'^#\s*(.+?)(?:\n|$)', content, re.MULTILINE)
#             logger.info(f"title_match:{title_match}")
#             if title_match:
#                 subject = title_match.group(1).strip()[:50]
#                 logger.info(f"subject:{subject}")
#             else:
#                 subject = "您要求的邮件"
#         logger.info(f"={subject}:{content[:200]},\n{attachment_path}")
        # 4. 实际发送邮件
        # return self._execute_email_send(to_email, subject, content, attachment_path)
        
    #没使用
    # def _handle_email_subtask(self, task: SubTask, user_input: str,
    #                         context: Optional[Dict], task_results: Dict) -> str:
    #     """处理邮件发送子任务"""
    #     user_input = user_input + task.params.get('email', '') if task.params else ''
    #     for dep_id in (task.depends_on or []):
    #         if dep_id in task_results:
    #             dep_result = task_results[dep_id]
    #     user_input += dep_result
    #     return self._dispatch_send_mail(user_input,context)
        
    #direct用
    def _dispatch_send_mail(self, user_input: str,
                              context: Optional[Dict]) -> str:
        logger.info(f"==============开始准备发邮件==================")
        # 1. 从 task.params 或 user_input 提取邮箱地址
        email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'

        matches = re.findall(email_pattern, user_input)
        to_email = matches[0] if matches else ''

        if not to_email:
            return "❌ 未找到收件人邮箱地址"

        # 2. 从前置任务结果获取内容/附件
        attachment_path = None
        prompt = f"""基于用户输入和上下文信息，提取邮件的主题、内容和附件。

用户输入：{user_input}
上下文信息：{context}

指令：
1. 优先从上下文中提取现有内容（标题、正文、文件路径），不要重新生成
2. 如果上下文有标题和短文，直接使用：标题→邮件主题，短文→邮件内容
3. 如果上下文有完整文章/文档：提取文件名作为主题，文章内容作为正文，提取文件路径
4. 如果上下文信息不足，生成专业的邮件内容，但要保留已有内容
5. 多个上下文时，优先使用最近的一条

返回JSON格式（严格遵守）：
{{
    "subject": "邮件主题（字符串）",
    "body": "邮件正文内容（字符串）",
    "attachment_path": "文件路径或null（如果无附件）"
}}"""

        model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"
        response = self.model_caller(prompt, model=model, temperature=0.3)
        result = self._extract_json(response) if response else None
        logger.info(f"send_mail:{result}")

        if not result:
            return "❌ LLM 提取邮件信息失败"

        # 优先使用 context 中存储的生成文件路径
        if context and context.get("generated_docx"):
            attachment_path = context.get("generated_docx")
            logger.info(f"[Email] 使用上下文中的 Word 文档: {attachment_path}")
        elif isinstance(result.get("attachment_path"), str):
            # 查找文件路径 - 优先查找本地生成的文件
            # 先检查是否有生成的 .docx 文件
            docx_match = re.search(r'([^\s]+\.docx)', result.get("attachment_path"))
            if docx_match:
                potential_path = docx_match.group(1)
                if os.path.exists(potential_path):
                    attachment_path = potential_path
                    logger.info(f"[Email] 检测到生成的 Word 文档: {potential_path}")

            # 其次查找其他格式的文件
            if not attachment_path:
                path_match = re.search(r'(/[^\s]+\.(docx|pptx|pdf|md|txt))', result.get("attachment_path"))
                if path_match and os.path.exists(path_match.group(1)):
                    attachment_path = path_match.group(1)
        # 3. 构建邮件主题
        subject = result.get('subject', '').strip()
        if not subject:
            # 从正文提取第一行作为主题
            body_text = result.get('body', '').strip()
            title_match = re.search(r'^#?\s*(.+?)(?:\n|$)', body_text)
            subject = title_match.group(1).strip()[:50] if title_match else "您要求的邮件"
        logger.info(f"={subject}:{result.get('body','根据您的要求，发送该邮件。谢谢')[:200]},{attachment_path}")
        # 4. 实际发送邮件
        result = self._execute_email_send(to_email, subject, result.get("body","根据您的要求，发送该邮件。\n\n谢谢"), attachment_path)
        logger.info(f"==============邮件发送结束，\n{result}")

    def _execute_email_send(self, to_email: str, subject: str, body: str,
                            attachment_path: str = None) -> str:
        # """实际执行邮件发送"
        logger.info(f"调用server的send_email功能，{subject}:{body}-\n{attachment_path},")
        try:
            # 导入 server.py 中的 send_email 函数
            import sys
            server_path = str(Path(__file__).parent)
            if server_path not in sys.path:
                sys.path.insert(0, server_path)
            from server import send_email

            success = send_email(to_email, subject, body, attachment_path)
            if success:
                result = f"✅ 邮件已成功发送到 {to_email}"
                if attachment_path:
                    result += f"（附件: {os.path.basename(attachment_path)}）"
                logger.info(f"[Email] {result}")
                return result
            else:
                logger.warning(f"[Email] 发送失败: {to_email}")
                return f"❌ 邮件发送失败"
        except Exception as e:
            logger.error(f"[Email] 发送异常: {e}")
            return f"❌ 邮件发送异常: {e}"
# 对于 skill 类型的任务（如 docx），当前实现只是通过 LLM 生成指令，而不是实际执行 skill。我需要修改 _execute_subtask_skill 方法，对 docx skill               
#   类型进行特殊处理，使其能够实际生成 Word 文件。
#与  _create_docx重复 todo
    # def _handle_docx_skill(self, task: SubTask, user_input: str,
    #                        context: Optional[Dict], task_results: Dict) -> str:
    #     """处理 docx skill - 实际生成 Word 文档"""
    #     logger.info(f"[Skill] 处理 docx skill: {task.name}")

    #     try:
    #         from docx import Document
    #         from docx.shared import Pt, RGBColor
    #         from docx.enum.text import WD_ALIGN_PARAGRAPH
    #         import datetime

    #         # 从前置任务结果提取内容
    #         content = ""
    #         title = "文档"

    #         for dep_id in (task.depends_on or []):
    #             if dep_id in task_results:
    #                 content = str(task_results[dep_id])
    #                 # 尝试从内容中提取标题
    #                 title_match = re.search(r'^#\s*(.+?)(?:\n|$)', content, re.MULTILINE)
    #                 if title_match:
    #                     title = title_match.group(1).strip()[:50]
    #                 break

    #         if not content:
    #             return "❌ 没有可用的内容来创建文档"

    #         # 创建新文档
    #         doc = Document()

    #         # 设置文档属性
    #         doc.core_properties.title = title
    #         doc.core_properties.author = "Claude AI"
    #         doc.core_properties.subject = task.name

    #         # 添加标题
    #         title_para = doc.add_paragraph(title, style='Heading 1')
    #         title_para_format = title_para.paragraph_format
    #         title_para_format.space_after = Pt(12)

    #         # 添加创建时间
    #         date_para = doc.add_paragraph(f"生成时间: {datetime.datetime.now().strftime('%Y年%m月%d日 %H:%M:%S')}")
    #         date_para.style = 'Normal'
    #         date_para.paragraph_format.space_after = Pt(12)

    #         # 添加内容（按行处理）
    #         for line in content.split('\n'):
    #             line = line.strip()
    #             if not line:
    #                 continue

    #             # 检测 markdown 标题
    #             if line.startswith('###'):
    #                 p = doc.add_paragraph(line[3:].strip(), style='Heading 3')
    #             elif line.startswith('##'):
    #                 p = doc.add_paragraph(line[2:].strip(), style='Heading 2')
    #             elif line.startswith('#'):
    #                 p = doc.add_paragraph(line[1:].strip(), style='Heading 1')
    #             elif line.startswith('- ') or line.startswith('* '):
    #                 # 列表项
    #                 p = doc.add_paragraph(line[2:].strip(), style='List Bullet')
    #             else:
    #                 # 普通段落
    #                 p = doc.add_paragraph(line)
    #                 p.paragraph_format.space_after = Pt(6)

    #         # 保存文档
    #         doc_filename = f"{title}.docx"
    #         doc.save(doc_filename)

    #         result = f"✅ Word 文档已成功生成: {doc_filename}"
    #         logger.info(f"[Skill] {result}")
    #         return result

    #     except Exception as e:
    #         logger.error(f"[Skill] docx 处理失败: {e}")
    #         return f"❌ Word 文档生成失败: {e}"

    def _build_task_chain_summary(self, tasks: List[SubTask], results: Dict) -> str:
        """构建任务链执行结果的总结"""
        lines = []
        lines.append("=" * 60)
        lines.append("📋 复杂任务链执行结果")
        lines.append("=" * 60)

        completed = sum(1 for t in tasks if t.status == "completed")
        failed = sum(1 for t in tasks if t.status in ("failed", "timeout"))
        blocked = sum(1 for t in tasks if t.status == "blocked")

        lines.append(f"\n📊 执行统计:")
        lines.append(f"  ✅ 成功: {completed}/{len(tasks)}")
        if failed > 0:
            lines.append(f"  ❌ 失败: {failed}")
        if blocked > 0:
            lines.append(f"  🚫 被阻止: {blocked}")

        lines.append(f"\n📝 详细步骤:")
        for task in tasks:
            status_emoji = {
                "completed": "✅",
                "failed": "❌",
                "timeout": "⏱️",
                "blocked": "🚫",
                "pending": "⏳",
                "running": "▶️"
            }.get(task.status, "❓")

            duration = ""
            if task.start_time and task.end_time:
                elapsed = (task.end_time - task.start_time).total_seconds()
                duration = f" ({elapsed:.1f}s)"

            lines.append(f"  {status_emoji} {task.task_id}: {task.name}{duration}")

            if task.error:
                lines.append(f"     错误: {task.error}")

            if task.status == "completed" and task.task_id in results:
                result_preview = str(results[task.task_id])
                if len(str(results[task.task_id])) > 100:
                    result_preview += "..."
                lines.append(f"     结果: {result_preview}")

        lines.append("\n" + "=" * 60)

        # 如果都成功，返回最后的结果
        if all(t.status == "completed" for t in tasks):
            lines.append("✨ 所有任务都已成功完成！")
            # 返回最后一个任务的结果作为最终输出
            last_task = tasks[-1]
            if last_task.task_id in results:
                lines.append(f"\n最终结果:\n{results[last_task.task_id]}")
        else:
            lines.append("⚠️ 任务链执行存在失败或被阻止的任务")

        return "\n".join(lines)

    def _dispatch_skill(self, intent: Intent, user_input: str, context: Optional[Dict]) -> str:
        """调度 Skill"""
        skill = self.discovery.get_skill_by_name(intent.target)
        logger.info(f"=========================该步需要执行的skill:{str(skill)}==============")
        if not skill:
            return f"❌ Skill '{intent.target}' 未找到"
                # 检测是否涉及文档生成，使用更长的超时
        # doc_keywords = ["转化为ppt", "转为ppt", "生成ppt", "创建ppt", "制作ppt",
        #                 "转化为docx", "转为docx", "生成docx", "创建docx",
        #                 "转化为word", "转为word", "生成word", "创建word",
        #                 "转化为文档", "转为文档", "生成文档"]
        # is_doc_gen = any(kw in user_input.lower() for kw in doc_keywords)
        # timeout = 900 if is_doc_gen else 300
        
        # 特殊处理 knowledge-base skill - 直接调用脚本
        if skill.name == "knowledge-base":
            return self._dispatch_knowledge_base(intent, user_input, context)

        # 构建 Skill 调用提示
        prompt = f"""请使用 {skill.name} Skill 来处理以下用户请求。

用户请求: {user_input}

Skill 描述: {skill.description}

【重要约束】
1. 只完成当前 Skill 的任务，不要包含其他步骤的代码（如发送邮件、其他文件操作等）
2. 如果生成代码，请确保代码可以独立执行
3. 生成的文件保存到当前目录
"""

        if intent.operation and intent.operation in skill.operations:
            op = skill.operations[intent.operation]
            prompt += f"\n建议操作: {intent.operation}\n原因: {'; '.join(op.semantic_triggers[:2])}"
        logger.info(f"==该步需要执行的skill promp:{str(prompt)}")
        # 从context中获取用户选择的模型，如果没有则使用默认值
        model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"

        # 调用 Claude 来执行 skill
        result = self.model_caller(prompt, model=model, temperature=0.3)
        logger.info("=================skill执行结果:\n{result}")
        # 特殊处理 docx skill - 如果输出包含 Python 代码，执行它来生成文件
        if skill.name == "docx" and "```python" in result:
            result = self._execute_docx_code(result, context)
            logger.info("=================生成docx:\n{result}")
            return result

        return result
    #执行代码
    def _execute_docx_code(self, llm_output: str, context: Optional[Dict]) -> str:
        """执行 docx skill 生成的 Python 代码，创建 Word 文档"""
        import re
        import tempfile
        import subprocess
        import os
        import time

        logger.info("[DOCX] 检测到 Python 代码，尝试执行生成 Word 文档...")

        # 提取 Python 代码块
        code_pattern = r'```python\s*(.*?)```'
        matches = re.findall(code_pattern, llm_output, re.DOTALL)

        if not matches:
            logger.warning("[DOCX] 未找到有效的 Python 代码块")
            return llm_output

        # 取最长的代码块（通常是主要代码）
        code = max(matches, key=len)

        # 移除发送邮件相关的代码（如果有的话）
        code = re.sub(r'def send_email.*?(?=\ndef |\nif __name__|$)', '', code, flags=re.DOTALL)
        code = re.sub(r'send_email.*?\(.*?\)', '# send_email removed', code)

        # 确保代码调用了 main 函数或 create_word_document 函数
        if 'if __name__' not in code:
            if 'def main' in code:
                code += '\n\nif __name__ == "__main__":\n    main()\n'
            elif 'def create_word_document' in code:
                code += '\n\nif __name__ == "__main__":\n    create_word_document()\n'

        # 创建临时文件执行代码
        try:
            with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False, encoding='utf-8') as f:
                f.write(code)
                temp_script = f.name

            logger.info(f"[DOCX] 执行临时脚本: {temp_script}")

            # 执行代码
            result = subprocess.run(
                ['python3', temp_script],
                capture_output=True,
                text=True,
                timeout=120,
                cwd=os.getcwd()
            )
            logger.info(f"===========运行生成docx代码的结果：{result}")
            # 清理临时文件
            os.unlink(temp_script)

            if result.returncode == 0:
                # 查找生成的 .docx 文件
                docx_files = [f for f in os.listdir('.') if f.endswith('.docx') and os.path.getmtime(f) > (time.time() - 60)]
                logger.info(f"========doc路径：{docx_files}")
                if docx_files:
                    newest_file = max(docx_files, key=os.path.getmtime)
                    file_path = os.path.abspath(newest_file)
                    logger.info(f"[DOCX] ✅ Word 文档已生成: {file_path}")

                    # 将文件路径存储到 context 中，供后续步骤使用
                    if context:
                        context['generated_docx'] = file_path
                    logger.info("=======文件路径: {file_path}\n\n{result.stdout},\n{str(context)}")

                    return f"✅ Word 文档已成功生成！\n\n文件路径: {file_path}\n\n{result.stdout}"
                else:
                    logger.warning("[DOCX] 代码执行成功但未找到生成的 .docx 文件")
                    return f"⚠️ 代码执行成功但未找到生成的文件\n\n输出: {result.stdout}"
            else:
                logger.error(f"[DOCX] 代码执行失败: {result.stderr}")
                return f"❌ 代码执行失败\n\n错误: {result.stderr}\n\n原始输出:\n{llm_output[:2000]}"

        except subprocess.TimeoutExpired:
            logger.error("[DOCX] 代码执行超时")
            return "❌ 代码执行超时（120秒）"
        except Exception as e:
            logger.error(f"[DOCX] 执行异常: {e}")
            return f"❌ 执行异常: {str(e)}\n\n原始输出:\n{llm_output[:2000]}"


    def _dispatch_agent(self, intent: Intent, user_input: str, context: Optional[Dict]) -> str:
        """调度 Agent"""
        agent = self.discovery.get_agent_by_name(intent.target)
        logger.info(f"==该步需要执行的skill promp:{str(agent)}")
        if not agent:
            return f"❌ Agent '{intent.target}' 未找到"

        # 构建 Agent 调用提示
        prompt = f"""请使用 {agent.name} Agent 来处理以下用户请求。

Agent 描述: {agent.description}
适用场景: {'; '.join(agent.use_cases)}

用户请求: {user_input}

请根据Agent类型的特点，使用相应的工具（如 Glob, Grep, Read, Task 等）来完成这个任务。
"""
        logger.info(f"==该步需要执行的skill promp:{str(prompt)}")
        # 从context中获取用户选择的模型，如果没有则使用默认值
        model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"

        return self.model_caller(prompt, model=model, temperature=0.3)

    def _dispatch_knowledge_base(self, intent: Intent, user_input: str, context: Optional[Dict]) -> str:
        """调度 knowledge-base skill - 直接调用Python脚本

        知识库选择优先级：
        1. 用户输入中显式指定（如"在国学知识库中"、"KBGX"） - 最高
        2. UI上的知识库选择 - 中等
        3. 自动检测 - 最低
        """
        import subprocess
        import re

        # 网页线程有输入后才有模型参数-从 context 中获取前端模型，如果没有则使用默认值
        model = normalize_model(context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1")
        logger.info(f"==前端输入模型：{model}, intent:{str(intent)}")
        # ========== 处理对话历史上下文 ==========
        history = context.get('history', []) if context else []
        original_query = user_input  # 保存原始用户输入
        previous_assistant_content = None  # 之前助手的回复内容
        previous_user_query = None  # 之前用户的查询

        if history:
            # 解析历史，提取有用信息
            for msg in reversed(history):  # 从最近的开始
                role = msg.get('role', '')
                content = msg.get('content', '')
                if role == 'assistant' and not previous_assistant_content:
                    # 获取最近的助手回复（可能是需要转换的内容）
                    previous_assistant_content = content
                elif role == 'user' and not previous_user_query:
                    # 获取最近的用户查询（可能是原始搜索词）
                    previous_user_query = content

            # 判断当前请求是否是引用上文的请求
            input_lower_check = user_input.lower()
            is_contextual_request = any(kw in input_lower_check for kw in [
                '上面', '上述', '刚才', '前面', '这个', '那个',
                '转换', '导出', '下载', '保存', '生成',
                '继续', '接着', 'word', 'docx', 'ppt', 'pptx'
            ])

            if is_contextual_request:
                logger.info(f"[KB] 检测到引用上文请求: {user_input}")

                # 如果是"继续搜索"类请求，使用之前的查询词
                if any(kw in input_lower_check for kw in ['继续', '接着', '再搜索', '重新搜索']):
                    if previous_user_query:
                        # 从之前的查询中提取关键词
                        original_query = previous_user_query
                        logger.info(f"[KB] 使用之前的查询: {original_query}")

                # 如果是转换请求（Word/PPT/MD），并且有之前的助手回复
                if any(kw in input_lower_check for kw in ['word', 'docx', 'ppt', 'pptx', '文档', '演示', 'md', 'markdown']):
                    if previous_assistant_content and len(previous_assistant_content) > 50:
                        # 检查是否是有效的内容（不是错误信息）
                        if not previous_assistant_content.startswith('Error') and not previous_assistant_content.startswith('❌'):
                            logger.info(f"[KB] 检测到转换请求，将使用之前的助手回复内容 (长度: {len(previous_assistant_content)})")
                            # 将之前的内容保存到context，供后续使用
                            context['content_to_convert'] = previous_assistant_content

                # ========== 直接处理 md 下载请求 ==========
                # 如果用户要求将上文内容保存为 md 下载
                if any(kw in input_lower_check for kw in ['md', 'markdown']) and any(kw in input_lower_check for kw in ['下载', '保存', '导出', '转换', '变成']):
                    if previous_assistant_content and len(previous_assistant_content) > 50:
                        if not previous_assistant_content.startswith('Error') and not previous_assistant_content.startswith('❌'):
                            logger.info(f"[KB] 直接处理 md 下载请求，使用之前的助手回复")
                            # 直接保存并返回下载链接
                            return self._save_and_return_download_link(previous_assistant_content, user_input, "md", context)

        input_lower = original_query.lower()

        # ========== 用 original_query 替换 user_input ==========
        # 确保后续所有操作使用正确的查询词，而不是"继续搜索"等引用上文的短语
        if original_query != user_input:
            logger.info(f"[KB] user_input 从 '{user_input[:50]}' 替换为 '{original_query[:50]}'")
            user_input = original_query

        # 优先级1：检查用户输入中是否显式指定了知识库
        kb_name = None
        explicit_kb_keywords = {
            "KBGX": ["kbgx", "国学知识库", "国学库", "国学"],
            "KBW": ["kbw", "微信知识库", "微信库", "微信"],
            "KB": ["kb", "通用知识库", "通用库"]
        }

        for kb, keywords in explicit_kb_keywords.items():
            for keyword in keywords:
                if keyword in input_lower:
                    kb_name = kb
                    logger.info(f"[KB] 从用户输入检测到显式KB指定: {kb_name}")
                    break
            if kb_name:
                break

        # 优先级2：如果用户输入中没有显式指定，使用UI选择
        if not kb_name and context and context.get("kb"):
            kb_name = str(context.get("kb")).upper()
            logger.info(f"[KB] 使用UI选择的知识库: {kb_name}")

        # 优先级3：如果都没有，从输入进行一般检测
        if not kb_name:
            kb_name = "KB"  # 默认
            # 再检查一遍关键词，以防UI也没有选择
            if "kbgx" in input_lower or "国学" in input_lower:
                kb_name = "KBGX"
            elif "kbw" in input_lower or "微信" in input_lower:
                kb_name = "KBW"
            logger.info(f"[KB] 自动检测的知识库: {kb_name}")

        kb_path = ""
        if self.kb_paths:
            kb_path = self.kb_paths.get(kb_name, "") or ""

        if not kb_path:
            return f"❌ 找不到知识库 {kb_name} 的路径"

        scripts_dir = Path("/home/will/.claude/skills/knowledge-base/scripts")

        # ✅ 关键修复：根据intent.operation立即决定操作，而非重新分析user_input
        if intent.operation == "file_download":
            # 下载现有的文件
            logger.info(f"[KB] 执行 file_download 操作（来自intent）")
            try:
                filename = intent.params.get("filename", "").strip()
                if not filename:
                    return "❌ 未指定文件名"

                logger.info(f"[KB] 开始搜索文件: {filename}")

                # 定义搜索顺序 - 与 server.py 中 _handle_download_file 一致
                search_dirs = [
                    *KB_PATHS.values(),  # 所有知识库目录
                    "/home/will/Downloads/opencode_p/tmp/",  # 新生成文件目录
                    "/home/will/Downloads/opencode_p/",  # 下载目录
                    "/home/will/Downloads",  # 用户下载目录
                ]

                file_path = None
                for search_dir in search_dirs:
                    candidate = os.path.join(search_dir, filename)
                    if os.path.exists(candidate) and os.path.isfile(candidate):
                        file_path = candidate
                        logger.info(f"[KB] 找到文件: {file_path}")
                        break

                if not file_path:
                    logger.warning(f"[KB] 文件未找到: {filename}")
                    return f"❌ 找不到文件: {filename}\n\n请确认文件名是否正确，或者文件是否存在于知识库中。"

                # 构建下载链接 - 使用完整URL，支持局域网访问
                from urllib.parse import quote
                encoded_filename = quote(filename)
                base_url = context.get("base_url", "") if context else ""
                download_url = f"{base_url}/api/download/{encoded_filename}"

                logger.info(f"[KB] 文件下载链接: {download_url}")

                return f"""✅ 文件已找到

📥 **下载文件**
[{filename}]({download_url})"""

            except Exception as e:
                logger.error(f"[KB] file_download 失败: {e}", exc_info=True)
                return f"❌ 文件下载失败: {e}"

        elif intent.operation == "list_docs":
            # 列出知识库文件
            logger.info(f"[KB] 执行 list_docs 操作（来自intent）")
            try:
                result = subprocess.run(
                    ["python3", str(scripts_dir / "list_documents.py"), "--kb-path", kb_name, "--db-path", "/home/will/Downloads/opencode_p/.knowledge_base"],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                return result.stdout if result.stdout else result.stderr
            except Exception as e:
                return f"❌ 列出文件失败: {e}"

        elif intent.operation == "context_download":
            # 从上下文（对话历史）获取内容并下载
            logger.info(f"[KB] 执行 context_download 操作（来自intent）")
            try:
                # 从对话历史中获取上一条助手回复
                previous_content = None
                if history:
                    for msg in reversed(history):
                        if msg.get('role') == 'assistant':
                            content = msg.get('content', '')
                            # 检查是否是有效内容（不是错误信息）
                            if content and len(content) > 50 and not content.startswith('❌') and not content.startswith('Error'):
                                previous_content = content
                                break

                if not previous_content:
                    return "❌ 未找到可下载的上文内容。请先进行查询或总结，然后再请求下载。"

                # 检测用户请求的格式
                output_format = intent.params.get('format', 'md')
                input_lower_check = user_input.lower()

                if any(kw in input_lower_check for kw in ['word', 'doc', 'docx']):
                    output_format = 'docx'
                elif any(kw in input_lower_check for kw in ['ppt', 'pptx', '演示', '幻灯']):
                    output_format = 'pptx'
                elif any(kw in input_lower_check for kw in ['md', 'markdown']):
                    output_format = 'md'

                logger.info(f"[KB] 从上下文下载，格式: {output_format}, 内容长度: {len(previous_content)}")

                # 保存并返回下载链接
                return self._save_and_return_download_link(previous_content, user_input, output_format, context)

            except Exception as e:
                return f"❌ 下载失败: {e}"

        elif intent.operation == "show_full_doc":
            # 显示文档全文
            logger.info(f"[KB] 执行 show_full_doc 操作（来自intent）")
            try:
                # 从intent参数中获取文档名称
                doc_name = intent.params.get('doc_name', '')
                if not doc_name:
                    # 从用户输入中提取
                    doc_name = user_input
                    for remove_word in ["显示", "查看", "看看", "看", "展示", "打开", "读",
                                        "阅读", "全文", "正文", "完整内容", "全部内容",
                                        "整篇", "的", "请", "帮我"]:
                        doc_name = doc_name.replace(remove_word, "")
                    doc_name = doc_name.strip()

                if not doc_name:
                    return "❌ 无法识别文档名称，请指定要查看的文档"

                # 在知识库索引中查找匹配的文档
                db_path = f"/home/will/Downloads/opencode_p/.knowledge_base/{kb_name.lower()}_index.json"

                if not os.path.exists(db_path):
                    return f"❌ 知识库索引文件不存在: {db_path}"

                with open(db_path, 'r', encoding='utf-8') as f:
                    index_data = json.load(f)

                docs = index_data.get('documents', {})
                matched_doc = None

                # 1. 精确匹配标题
                for doc_id, doc_info in docs.items():
                    title = doc_info.get('title', '')
                    if title == doc_name or doc_name == title:
                        matched_doc = (doc_id, doc_info)
                        break

                # 2. 标题包含匹配
                if not matched_doc:
                    for doc_id, doc_info in docs.items():
                        title = doc_info.get('title', '')
                        if doc_name in title or title in doc_name:
                            matched_doc = (doc_id, doc_info)
                            break

                # 3. 文件名包含匹配
                if not matched_doc:
                    for doc_id, doc_info in docs.items():
                        # 去掉扩展名比较
                        base_name = doc_id.rsplit('.', 1)[0] if '.' in doc_id else doc_id
                        if doc_name in base_name or base_name in doc_name:
                            matched_doc = (doc_id, doc_info)
                            break

                if not matched_doc:
                    # 列出可用文档帮助用户选择
                    doc_titles = [doc_info.get('title', doc_id) for doc_id, doc_info in list(docs.items())[:10]]
                    titles_str = '\n'.join(f"  - {t}" for t in doc_titles)
                    return f"❌ 在知识库 {kb_name} 中未找到名为 '{doc_name}' 的文档\n\n可用文档：\n{titles_str}"

                doc_id, doc_info = matched_doc
                title = doc_info.get('title', doc_name)

                # 优先从索引中获取已提取的内容
                content = doc_info.get('content', '')

                if not content:
                    # 尝试直接读取文件（仅对文本格式有效）
                    file_path = os.path.join(kb_path, doc_id)
                    if os.path.exists(file_path):
                        try:
                            with open(file_path, 'r', encoding='utf-8') as f:
                                content = f.read()
                        except Exception:
                            pass

                if content:
                    logger.info(f"[KB] 显示全文: {title}, 内容长度: {len(content)}")
                    return f"# {title}\n\n{content}"
                else:
                    return f"❌ 找到文档 '{title}' 但无法读取内容"

            except Exception as e:
                return f"❌ 显示全文失败: {e}"

        elif intent.operation == "summarize_doc":
            # 总结文档
            logger.info(f"[KB] 执行 summarize_doc 操作（来自intent）")
            try:
                # 提取文档名称
                doc_name = intent.params.get('doc_name', '')
                if not doc_name:
                    doc_name = user_input.replace("总结", "").replace("summarize", "").strip()

                # ========== 优先尝试直接从索引读取文档内容 ==========
                db_path = f"/home/will/Downloads/opencode_p/.knowledge_base/{kb_name.lower()}_index.json"

                if os.path.exists(db_path):
                    with open(db_path, 'r', encoding='utf-8') as f:
                        index_data = json.load(f)

                    docs = index_data.get('documents', {})
                    matched_doc = None

                    # 1. 精确匹配标题
                    for doc_id, doc_info in docs.items():
                        title = doc_info.get('title', '')
                        if title == doc_name or doc_name == title:
                            matched_doc = (doc_id, doc_info)
                            break

                    # 2. 标题包含匹配
                    if not matched_doc:
                        for doc_id, doc_info in docs.items():
                            title = doc_info.get('title', '')
                            if doc_name in title or title in doc_name:
                                matched_doc = (doc_id, doc_info)
                                break

                    # 3. 文件名包含匹配
                    if not matched_doc:
                        for doc_id, doc_info in docs.items():
                            base_name = doc_id.rsplit('.', 1)[0] if '.' in doc_id else doc_id
                            if doc_name in base_name or base_name in doc_name:
                                matched_doc = (doc_id, doc_info)
                                break

                    if matched_doc:
                        doc_id, doc_info = matched_doc
                        title = doc_info.get('title', doc_name)
                        content = doc_info.get('content', '')

                        if not content:
                            # 尝试直接读取文件
                            file_path = os.path.join(kb_path, doc_id)
                            if os.path.exists(file_path):
                                try:
                                    with open(file_path, 'r', encoding='utf-8') as f:
                                        content = f.read()
                                except Exception:
                                    pass

                        if content:
                            logger.info(f"[KB] 直接从索引获取文档内容: {title}, 长度: {len(content)}")

                            # 使用 LLM 进行总结
                            processed_result = self._llm_post_process(user_input, content, model=model)

                            # 检查是否需要下载
                            if "下载" in user_input or "保存" in user_input or "导出" in user_input:
                                return self._save_and_return_download_link(processed_result, user_input, "md", context)
                            return processed_result
                        else:
                            logger.warning(f"[KB] 找到文档但无法获取内容: {title}")

                # ========== 回退到搜索方式 ==========
                logger.info(f"[KB] 未能直接匹配文档，使用搜索方式")
                query = user_input.replace("总结", "").replace("summarize", "").strip()
                logger.info(f"[KB] 需要后处理，生成多个查询以收集完整信息")
                all_results = []

                # 根据原始查询生成多个搜索角度的查询语句（传递 model 参数以使用 LLM 扩展关键词）
                queries = self._generate_search_queries(user_input, query, model)
                logger.info(f"[KB] summarize_doc生成的查询语句: {queries}")

                # 使用线程池并行查询，提高效率
                with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
                    # 提交所有查询任务
                    future_to_query = {
                        executor.submit(self._search_kb_single, search_query, kb_name): search_query
                        for search_query in queries
                    }

                    # 收集结果
                    for future in concurrent.futures.as_completed(future_to_query):
                        search_query = future_to_query[future]
                        try:
                            result = future.result(timeout=120)
                            if result and "No relevant documents found" not in result:
                                all_results.append(result)
                                logger.info(f"[KB] 查询 '{search_query}' 返回结果:\n{result}")
                        except Exception as e:
                            logger.warning(f"[KB] 查询 '{search_query}' 失败: {e}")

                if all_results:
                    # 去重相似的结果，避免重复内容
                    unique_results = self._deduplicate_search_results(all_results, similarity_threshold=0.7)
                    combined_output = "\n\n---\n\n".join(unique_results)
                    logger.info(f"[KB] 合并 {len(unique_results)} 个去重结果，总长度: {len(combined_output)}")

                    processed_result = self._llm_post_process(user_input, combined_output, model=model)

                    if "下载" in user_input or "保存" in user_input or "导出" in user_input:
                        return self._save_and_return_download_link(processed_result, user_input, "md", context)
                    return processed_result
                else:
                    logger.warning(f"[KB] 多查询未返回任何结果")
                    return f"❌ 在知识库 {kb_name} 中未找到文档 '{doc_name}'，请检查文档名称是否正确"

            except Exception as e:
                return f"❌ 总结文档失败: {e}"

        elif intent.operation == "extract_chapter":
            # 提取章节
            logger.info(f"[KB] 执行 extract_chapter 操作（来自intent）")
            try:
                doc_name = intent.params.get('doc_name', '')
                chapter = intent.params.get('chapter', '')

                if not doc_name or not chapter:
                    return "❌ 需要指定文档名和章节名"

                result = subprocess.run(
                    ["python3", str(scripts_dir / "get_chapter.py"), doc_name, chapter, "--kb-path", kb_path],
                    capture_output=True,
                    text=True,
                    timeout=600
                )
                output = result.stdout if result.stdout else result.stderr

                if "下载" in user_input or "保存" in user_input:
                    return self._save_and_return_download_link(output, user_input, "md", context)
                return output
            except Exception as e:
                return f"❌ 提取章节失败: {e}"

        elif intent.operation == "compare_docs":
            # 对比文档
            logger.info(f"[KB] 执行 compare_docs 操作（来自intent）")
            try:
                doc1 = intent.params.get('doc1', '')
                doc2 = intent.params.get('doc2', '')

                if not doc1 or not doc2:
                    return "❌ 需要指定两个要对比的文档"

                # 构建对比查询
                query = f"对比{doc1}和{doc2}"
                result = subprocess.run(
                    ["python3", str(scripts_dir / "search_knowledge_base.py"), query,
                     "--kb-path", f"/home/will/Downloads/opencode_p/{kb_name}",
                     "--db-path", f"/home/will/Downloads/opencode_p/.knowledge_base/{kb_name.lower()}_index.json",
                     "--format", "answer"],
                    capture_output=True,
                    text=True,
                    timeout=120
                )
                output = result.stdout if result.stdout else result.stderr

                # 使用LLM进行对比分析
                if output and "No relevant documents found" not in output:
                    processed_result = self._llm_post_process(user_input, output, model=model)
                    if "下载" in user_input or "保存" in user_input:
                        return self._save_and_return_download_link(processed_result, user_input, "md", context)
                    return processed_result
                return output
            except Exception as e:
                return f"❌ 对比文档失败: {e}"

        elif intent.operation == "index_docs":
            # 索引文档
            logger.info(f"[KB] 执行 index_docs 操作（来自intent）")
            try:
                result = subprocess.run(
                    ["python3", str(scripts_dir / "build_index.py"), "--kb-path", kb_name],
                    capture_output=True,
                    text=True,
                    timeout=120
                )
                output = result.stdout if result.stdout else result.stderr
                return f"✅ 文档索引完成\n\n{output}"
            except Exception as e:
                return f"❌ 索引文档失败: {e}"

        elif intent.operation == "document_generation":
            # 文档生成 - 将用户提供的内容转换为指定格式
            logger.info(f"[KB] 执行 document_generation 操作")
            try:
                # os, datetime, url_quote 已在文件顶部导入

                # 获取格式
                output_format = intent.params.get("format", "pptx") if intent.params else "pptx"
                format_mapping = {"ppt": "pptx", "word": "docx", "doc": "docx"}
                output_format = format_mapping.get(output_format, output_format)

                # 优先从context中获取要转换的内容（来自引用上文的请求）
                # 如果context中有保存的content_to_convert，使用它；否则从user_input中提取
                if context and context.get('content_to_convert'):
                    content = context.get('content_to_convert')
                    logger.info(f"[KB] 使用context中保存的内容，长度: {len(content)}")
                else:
                    # 提取要转换的内容（去除转换指令部分）
                    content = user_input
                    remove_patterns = [
                        r"^将下列内容转化为\w+[：:]?\s*",
                        r"^将以下内容转化为\w+[：:]?\s*",
                        r"^把下列内容转为\w+[：:]?\s*",
                        r"^将.*?转化为\w+[：:]?\s*",
                        r"^转化为\w+[：:]?\s*",
                        r"^转为\w+[：:]?\s*",
                        r"^生成\w+[：:]?\s*",
                        r"^创建\w+[：:]?\s*",
                    ]
                    for pattern in remove_patterns:
                        content = re.sub(pattern, "", content, flags=re.IGNORECASE)
                    content = content.strip()

                if not content:
                    return "❌ 未检测到要转换的内容，请提供要转换的文本"

                # 生成文件名
                # 尝试从内容中提取标题
                title_match = re.search(r'^#\s*(.+?)(?:\n|$)', content)
                if title_match:
                    title = title_match.group(1).strip()
                    safe_title = re.sub(r'[<>:"/\\|?*]', '', title)[:30]
                else:
                    safe_title = "文档"

                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                save_dir = "/home/will/Downloads/opencode_p/tmp/"
                os.makedirs(save_dir, exist_ok=True)

                if output_format == "pptx":
                    filename = f"{safe_title}_{timestamp}.pptx"
                    file_path = os.path.join(save_dir, filename)
                    logger.info(f"[KB] 调用 _create_pptx 创建演示文稿: {file_path}")
                    self._create_pptx(content, file_path)
                    #todo
                    result = self._dispatch_skill('pptx', "将此内容做成专业的ppt文档，根据文档内容自动选择合适的模版、颜色", content) #用skill输入和输出不同
                    logger.info(f"pptx result:{result}")
                elif output_format == "docx":
                    filename = f"{safe_title}_{timestamp}.docx"
                    file_path = os.path.join(save_dir, filename)
                    logger.info(f"[KB] 调用 _create_docx 创建文档: {file_path}")
                    self._create_docx(content, file_path)
                    #todo
                    result = self._dispatch_skill('docx', "将此内容做成专业的word文档，颜色使用专业深蓝色", content) #用skill输入和输出不同
                    logger.info(f"docx result:{result}")
                else:
                    # 默认保存为 markdown
                    filename = f"{safe_title}_{timestamp}.md"
                    file_path = os.path.join(save_dir, filename)
                    with open(file_path, 'w', encoding='utf-8') as f:
                        f.write(content)

                # 检查文件是否创建成功
                if os.path.exists(file_path):
                    file_size = os.path.getsize(file_path)
                    encoded_filename = url_quote(filename)
                    base_url = context.get("base_url", "") if context else ""
                    download_url = f"{base_url}/api/download/{encoded_filename}"
                    format_names = {"pptx": "PowerPoint演示文稿", "docx": "Word文档", "md": "Markdown"}
                    format_name = format_names.get(output_format, output_format.upper())

                    return f"""✅ {format_name}已创建

📥 **下载{format_name}**
[{format_name}: {filename}]{download_url}"""
                else:
                    return f"❌ 文件创建失败，请重试"

            except Exception as e:
                logger.error(f"[KB] document_generation 失败: {e}", exc_info=True)
                return f"❌ 文档生成失败: {e}"
        #发送邮件
        elif intent.operation == "email_send":
            # 邮件发送 - 发送内容或文件到指定邮箱
            logger.info(f"[KB] 执行 email_send 操作")
            try:
                email_address = intent.params.get("email", "") if intent.params else ""
                output_format = intent.params.get("format", "text") if intent.params else "text"
                kb = intent.params.get("kb", "KB") if intent.params else "KB"

                if not email_address:
                    return "❌ 未检测到收件人邮箱地址"

                logger.info(f"[KB] 准备发送邮件到: {email_address}, 格式: {output_format}")

                # 获取要发送的内容
                if context and context.get('content_to_convert'):
                    content = context.get('content_to_convert')
                    logger.info(f"[KB] 使用context中保存的内容，长度: {len(content)}")
                else:
                    # 从user_input中提取内容
                    content = user_input
                    # 移除邮箱地址和邮件相关的关键词
                    import re as re_module
                    email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
                    content = re_module.sub(email_pattern, '', content)
                    content = content.replace('发送到', '').replace('发到', '').replace('邮箱', '').replace('邮件', '')
                    content = content.strip()

                if not content:
                    return "❌ 未检测到要发送的内容"

                # 构建邮件主题和内容
                title_match = re.search(r'^#\s*(.+?)(?:\n|$)', content)
                if title_match:
                    subject = title_match.group(1).strip()[:50]
                else:
                    subject = "来自Web服务器的内容"

                # 调用 server.py 中的邮件发送接口
                try:
                    import sys
                    import os as os_module
                    sys.path.insert(0, str(Path(__file__).parent.parent / 'scripts'))

                    # 这里我们通过 HTTP 调用 server.py 的邮件接口
                    # 或者直接导入并调用邮件函数
                    # 为了简化，我们返回一个提示，让 server.py 处理邮件发送

                    return f"""📧 邮件发送信息:
- 收件人: {email_address}
- 主题: {subject}
- 格式: {output_format}
- KB: {kb}

请通过 API 调用 /api/email/send-content 来完成邮件发送。

邮件内容已准备就绪。"""

                except Exception as e:
                    logger.error(f"[KB] 邮件发送失败: {e}")
                    return f"❌ 邮件发送失败: {e}"

            except Exception as e:
                logger.error(f"[KB] email_send 失败: {e}", exc_info=True)
                return f"❌ 邮件操作失败: {e}"

        elif intent.operation == "search":

            # ========== 优先检查：是否是显示特定文档内容的请求 ==========
            # 如果用户输入包含"显示/查看/看" + "内容"，先尝试精确匹配文档名
            display_keywords = ["显示", "查看", "看看", "看", "展示", "打开", "读", "阅读"]
            content_keywords = ["内容", "全文", "正文"]
            has_display_keyword = any(kw in user_input for kw in display_keywords)
            has_content_keyword = any(kw in user_input for kw in content_keywords)

            if has_display_keyword and has_content_keyword:
                # 提取可能的文档名
                potential_doc_name = user_input
                for remove_word in display_keywords + content_keywords + ["的", "请", "帮我", "我想", "我要"]:
                    potential_doc_name = potential_doc_name.replace(remove_word, "")
                potential_doc_name = potential_doc_name.strip()

                if potential_doc_name and len(potential_doc_name) >= 2:
                    # 尝试在知识库索引中查找匹配的文档
                    db_path = f"/home/will/Downloads/opencode_p/.knowledge_base/{kb_name.lower()}_index.json"

                    if os.path.exists(db_path):
                        try:
                            with open(db_path, 'r', encoding='utf-8') as f:
                                index_data = json.load(f)

                            docs = index_data.get('documents', {})
                            matched_doc = None

                            # 1. 精确匹配标题或文件名（去掉扩展名）
                            for doc_id, doc_info in docs.items():
                                title = doc_info.get('title', '')
                                base_name = doc_id.rsplit('.', 1)[0] if '.' in doc_id else doc_id

                                if (potential_doc_name == title or
                                    potential_doc_name == base_name or
                                    potential_doc_name == doc_id):
                                    matched_doc = (doc_id, doc_info)
                                    logger.info(f"[KB] 精确匹配到文档: {doc_id}")
                                    break

                            # 2. 标题或文件名包含匹配
                            if not matched_doc:
                                for doc_id, doc_info in docs.items():
                                    title = doc_info.get('title', '')
                                    base_name = doc_id.rsplit('.', 1)[0] if '.' in doc_id else doc_id

                                    if (potential_doc_name in title or title in potential_doc_name or
                                        potential_doc_name in base_name or base_name in potential_doc_name):
                                        matched_doc = (doc_id, doc_info)
                                        logger.info(f"[KB] 部分匹配到文档: {doc_id}")
                                        break

                            # 如果找到匹配的文档，直接读取并返回内容
                            if matched_doc:
                                doc_id, doc_info = matched_doc
                                file_path = doc_info.get('file_path', '')

                                if file_path and os.path.exists(file_path):
                                    title = doc_info.get('title', doc_id)
                                    with open(file_path, 'r', encoding='utf-8') as f:
                                        content = f.read()
                                    logger.info(f"[KB] 显示文档内容: {title}, 内容长度: {len(content)}")
                                    return f"# {title}\n\n{content}"
                                else:
                                    logger.warning(f"[KB] 文档路径不存在: {file_path}")

                        except Exception as e:
                            logger.error(f"[KB] 查找文档失败: {e}")

            # ========== 继续原有的搜索逻辑 ==========
            download_keywords = ["下载", "保存", "导出", "生成文档", "形成文档", "存为"]
            needs_download = any(kw in user_input for kw in download_keywords)

            process_keywords = ["总结", "汇总", "精华", "核心", "对比", "比较", "分析", "说明", "解释"]
            needs_processing = any(kw in user_input for kw in process_keywords)

            format_keywords = {
                "docx": ["word", "docx", "文档", "word文档", "WORD"],
                "pptx": ["ppt", "pptx", "演示", "幻灯", "PPT", "演示文稿"]
            }

            detected_formats = []
            for fmt, keywords in format_keywords.items():
                if any(kw in user_input.lower() for kw in keywords):
                    detected_formats.append(fmt)

            output_format = detected_formats[0] if detected_formats else "md"

            if "word" in user_input.lower() or "文档" in user_input:
                output_format = "docx"
            elif "ppt" in user_input.lower() or "幻灯" in user_input or "演示" in user_input:
                output_format = "pptx"

            logger.info(f"[KB] 检测到输出格式: {output_format}")

            query = user_input.replace("搜索", "").replace("查找", "").replace("问", "")
            query = user_input.replace("搜索", "").replace("查找", "").replace("问", "")
            query = query.replace(f"{kb_name}知识库", "").replace("国学知识库", "").strip()

            # 移除引号
            query = query.replace('"', '').replace("'", "").replace("'", "").strip()

            if not query:
                query = user_input # 如果提取失败使用原始输入

            try:
                # ✅ 新增：如果需要后处理，生成多个查询关键词以收集足够的信息
                if needs_processing:
                    logger.info(f"[KB] 需要后处理，生成多个查询以收集完整信息")
                    all_results = []

                    # 根据原始查询生成多个搜索角度的查询语句（传递 model 参数以使用 LLM 扩展关键词）
                    queries = self._generate_search_queries(user_input, query, model)
                    logger.info(f"[KB] search生成的查询语句: {queries}")

                    # 使用线程池并行查询，提高效率
                    with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
                        # 提交所有查询任务
                        future_to_query = {
                            executor.submit(self._search_kb_single, search_query, kb_name): search_query
                            for search_query in queries
                        }

                        # 收集结果
                        for future in concurrent.futures.as_completed(future_to_query):
                            search_query = future_to_query[future]
                            try:
                                result = future.result(timeout=120)
                                if result and "No relevant documents found" not in result:
                                    all_results.append(result)
                                    logger.info(f"[KB] 查询 '{search_query}' 返回结果:\n{result}")
                            except Exception as e:
                                logger.warning(f"[KB] 查询 '{search_query}' 失败: {e}")

                    if all_results:
                        # 去重相似的结果，避免重复内容
                        unique_results = self._deduplicate_search_results(all_results, similarity_threshold=0.7)
                        combined_output = "\n\n---\n\n".join(unique_results)
                        logger.info(f"[KB] 合并 {len(unique_results)} 个去重结果，总长度: {len(combined_output)}")

                        # 文档生成需要更长的超时时间（15分钟）
                        post_timeout = 900 if needs_download else 300
                        processed_result = self._llm_post_process(user_input, combined_output, timeout=post_timeout, model=model)

                        if needs_download:
                            return self._save_and_return_download_link(processed_result, user_input, output_format, context)
                        return processed_result
                    else:
                        logger.warning(f"[KB] 多查询未返回任何结果，使用单查询")

                # 单查询搜索（如果不需要后处理或多查询失败）
                output = self._search_kb_single(query, kb_name)

                # 如果需要处理但前面没有处理，则进行LLM后处理
                if needs_processing and output and "No relevant documents found" not in output:
                    # 文档生成需要更长的超时时间（15分钟）
                    post_timeout = 900 if needs_download else 300
                    processed_result = self._llm_post_process(user_input, output, timeout=post_timeout, model=model)

                    if needs_download:
                        return self._save_and_return_download_link(processed_result, user_input, output_format, context)
                    return processed_result

                if needs_download:
                    return self._save_and_return_download_link(output, user_input, output_format, context)

                return output

            except Exception as e:
                return f"❌ 搜索失败: {e}"
# 用户输入 "总结国学知识库中道德经的精华并下载"                                                                                                                                  
#       ↓                                                                                                                                                                          
#   1. 提取主体：道德经                                                                                                                                                            
#       ↓                                                                                                                                                                          
#   2. LLM扩展关键词 → 道无为, 德治天下, 法自然                                                                                                                                    
#       ↓                                                                                                                                                                          
#   3. 生成多个查询 → [道德经, 道德经道无为, 道德经德治天下, 道德经法自然, ...]                                                                                                    
#       ↓                                                                                                                                                                          
#   4. 并行执行4个查询（ThreadPoolExecutor）                                                                                                                                       
#       ↓                                                                                                                                                                          
#   5. 收集结果 → 去重相似内容                                                                                                                                                     
#       ↓                                                                                                                                                                          
#   6. LLM后处理：生成至少10条核心智慧要点                                                                                                                                         
#       ↓                                                                                                                                                                          
#   7. 保存文件并返回下载链接                                                             
    def _expand_subject_keywords(self, core_entity: str, model: str = "nvidia/minimaxai/minimax-m2.1") -> list:
        """使用 LLM 扩展主体相关的专业词汇

        Args:
            core_entity: 核心主体，如"道德经"
            model: 使用的模型

        Returns:
            相关专业词汇列表
        """
        try:
            prompt = f"""请为"{core_entity}"列出3-5个与其核心内容高度相关的专业词汇或关键词。
要求：
1. 每个词汇2-6个字
2. 这些词汇应该是在学术性文本、经典注释中频繁出现的术语
3. 只返回关键词，用逗号分隔，不要解释

示例：
- 输入：论语
- 输出：仁义礼乐,君子之道,修身齐家,孔孟之道

请为"{core_entity}"提供关键词："""

            result = self.model_caller(prompt, model=model, temperature=0.3)

            # 清理并提取关键词
            if result:
                # 移除可能的标点和空白
                cleaned = result.strip().replace('、', ',').replace('，', ',').replace('并', ',')
                keywords = [k.strip() for k in cleaned.split(',') if k.strip()]
                logger.info(f"[KB] LLM扩展关键词: {keywords}")
                return keywords[:5]  # 最多返回5个

        except Exception as e:
            logger.warning(f"[KB] LLM扩展关键词失败: {e}")

        return []

    def _generate_search_queries(self, original_input: str, main_query: str, model: str = "nvidia/minimaxai/minimax-m2.1") -> list:
        """根据原始查询生成多个搜索关键词（加强版 V3）

        用于在需要总结/对比时收集足够的相关信息，使用 LLM 智能扩展关键词

        Args:
            original_input: 用户原始输入
            main_query: 提取后的主查询
            model: 使用的 LLM 模型
        """
        # ========== 检测是否是长上下文，避免错误扩展 ==========
        # 如果输入包含历史对话标记或过长，不进行扩展
        is_long_context = (
            len(original_input) > 200 or
            "对话历史" in original_input or
            "[用户]:" in original_input or
            "[助手]:" in original_input or
            "当前用户请求" in original_input
        )

        if is_long_context:
            logger.warning(f"[KB] 检测到长上下文输入 (长度: {len(original_input)})，跳过查询扩展")
            # 只使用主查询，不扩展
            clean_query = re.sub(r'(总结|精华|下载|分析|的|并|知识库|KB|KBW|KBGX|中)', '', main_query).strip()
            return [main_query, clean_query] if clean_query and clean_query != main_query else [main_query]

        # 1. 基础清理
        clean_query = re.sub(r'(总结|精华|下载|分析|的|并|知识库|KB|KBW|KBGX|中)', '', main_query).strip()

        # 2. 识别核心主体（例如：道德经）
        core_entity = clean_query if clean_query else main_query

        queries = [main_query]  # 首先添加主查询

        # 3. 注入领域专业词汇（使用 LLM 扩展）
        domain_extras = self._expand_subject_keywords(core_entity, model)
        if domain_extras:
            # 随机挑选或全部加入，增强 BM25 命中率
            for i in domain_extras:
                queries.append(f"{core_entity} {''.join(i)}")

        # 4. 根据意图组合（精华/总结）
        if "精华" in original_input or "总结" in original_input or "核心" in original_input:
            # 不要只加"精华"，要加专业术语
            queries.extend([
                # f"{core_entity} 核心旨归",  # 专业词汇更容易命中高质量学术论文
                # f"{core_entity} 思想体系",
                f"{core_entity} 经典语录",
                f"{core_entity} 释义"
            ])

        # 5. 特殊处理"下载" —— 转化为寻找"全文"或"资源"
        # if "下载" in original_input:
        #     queries.extend([
        #         f"{core_entity} 全文",
        #         f"{core_entity} 整理本",
        #         f"{core_entity} 完整版"
        #     ])

        elif "对比" in original_input or "比较" in original_input:
            # 为对比查询，提取两个对象
            parts = main_query.split("和")
            if len(parts) >= 2:
                queries.append(parts[0].strip())
                queries.append(parts[1].strip())
                queries.append(f"{parts[0].strip()} 特点")
                queries.append(f"{parts[1].strip()} 特点")

        elif "说明" in original_input or "解释" in original_input or "分析" in original_input:
            # 为说明/解释/分析查询，生成多维度搜索
            queries.extend([
                f"{main_query} 背景",
                f"{main_query} 意义",
                f"{main_query} 应用",
                f"{main_query} 分析"
            ])

        # 移除重复的查询
        queries = list(dict.fromkeys(queries))
        logger.info(f"[KB] 最终查询列表: {queries}")
        return queries

    def _deduplicate_search_results(self, results: list, similarity_threshold: float = 0.7) -> list:
        """去重搜索结果，过滤掉内容过于相似的结果

        Args:
            results: 搜索结果列表，每个元素是一个字符串
            similarity_threshold: 相似度阈值，超过此值则视为重复

        Returns:
            去重后的结果列表
        """
        import hashlib

        def get_text_hash(text: str, length: int = 100) -> str:
            """计算文本的简短哈希，用于快速比较"""
            # 取文本开头和结尾的组合作为指纹
            short_text = text[:length//2] + text[-length//2:] if len(text) > length else text
            return hashlib.md5(short_text.strip().encode()).hexdigest()

        seen_hashes = {}  # hash -> original text
        unique_results = []

        for result in results:
            if not result or not result.strip():
                continue

            # 方法1：基于内容指纹的快速去重
            content_hash = get_text_hash(result)

            if content_hash not in seen_hashes:
                seen_hashes[content_hash] = result
                unique_results.append(result)
            else:
                logger.info(f"[KB] 发现重复结果，已过滤")

            # 方法2：检查是否与已有结果内容过于相似（字符串相似度）
            # 只针对较短的结果进行详细检查，避免性能问题
            if len(result) < 500:
                for existing in unique_results:
                    # 简单的相似度检查：看多少字符是相同的
                    common_chars = sum(1 for c in result if c in existing)
                    similarity = common_chars / min(len(result), len(existing))

                    if similarity > similarity_threshold:
                        logger.info(f"[KB] 发现相似结果 (相似度={similarity:.2f})，已过滤")
                        if result in unique_results:
                            unique_results.remove(result)
                        break

        logger.info(f"[KB] 去重: {len(results)} → {len(unique_results)} 个结果")
        return unique_results

    def _search_kb_single(self, query: str, kb_name: str) -> str:
        """执行单一的知识库搜索"""
        import subprocess
        from pathlib import Path

        scripts_dir = Path("/home/will/.claude/skills/knowledge-base/scripts")
        kb_full_path = f"/home/will/Downloads/opencode_p/{kb_name}"
        db_path = f"/home/will/Downloads/opencode_p/.knowledge_base/{kb_name.lower()}_index.json"

        try:
            result = subprocess.run(
                ["python3", str(scripts_dir / "search_knowledge_base.py"), query,
                 "--kb-path", kb_full_path, "--db-path", db_path, "--format", "answer"],
                capture_output=True,
                text=True,
                timeout=180
            )
            return result.stdout if result.stdout else result.stderr
        except Exception as e:
            logger.warning(f"[KB] 单查询 '{query}' 执行失败: {e}")
            return ""
#   1. 修改 _llm_post_process 函数签名（第 1650 行）                                                                                                                                                 
#     - 添加 model: str = None 参数                                                                                                                                                                  
#   2. 在调用 self.model_caller 时传递 model 参数（第 1741-1744 行）                                                                                                                                 
#     - 如果提供了 model，则显式传递它                                                                                                                                                               
#   3. 更新所有调用 _llm_post_process 的位置（共 5 处）                                                                                                                                              
#     - 第 1133 行（summarize_doc）                                                                                                                                                                  
#     - 第 1177 行（summarize_doc 多查询）                                                                                                                                                           
#     - 第 1238 行（compare_docs）                                                                                                                                                                   
#     - 第 1414 行（search_kb 多查询）                                                                                                                                                               
#     - 第 1429 行（search_kb 单查询）                                                                                                                                                               
                                           
    def _llm_post_process(self, user_input: str, search_result: str, timeout: int = 300, model: str = None) -> str:
        """使用LLM对搜索结果进行后处理（总结、对比、结构化等）

        Args:
            user_input: 用户输入
            search_result: 搜索结果
            timeout: 超时时间（秒），文档生成时使用更长的超时
            model: 模型名称，如果为None则使用默认值
        """
        import logging
        logger = logging.getLogger(__name__)

        # 提取主体（如"道德经"）
        def extract_subject(text: str) -> str:
            """从用户输入中提取主体"""
            # 移除常见的动词和修饰词
            subject = text
            for keyword in ["总结", "精华", "核心", "下载", "分析", "说明", "解释", "国学知识库中", "知识库中", "中"]:
                subject = subject.replace(keyword, "")
            subject = subject.strip()
            return subject if subject else "相关内容"

        subject = extract_subject(user_input)

        # 根据请求类型构建不同的处理提示
        if "对比" in user_input or "比较" in user_input:
            prompt = f"""请基于以下来自知识库的内容，进行详细的比较和对比分析：

【知识库搜索结果】
{search_result}

【分析要求】
1. 识别涉及的主要对象和观点
2. 进行系统的对比分析（相同点、不同点、优劣分析）
3. 用表格或清晰的列表形式呈现对比结果
4. 分析造成差异的原因
5. 给出结论性评价

请用结构化的Markdown格式组织答案，确保逻辑清晰、数据准确。"""

        elif "精华" in user_input or "核心" in user_input or "总结" in user_input:
            prompt = f"""请总结{subject}的核心智慧要点：

【知识库搜索结果】
{search_result}

【总结要求】
1. 请提取至少10条核心智慧要点
2. 每条用简练的语言概括主要思想
3. 重点涵盖该主题的重要思想观点
4. 对每个要点进行简要解释说明（2-3句话）
5. 按重要性或逻辑顺序排列
6. 用清晰的列表格式呈现
7. 保留重要的原文引用和出处

请用精炼、准确的中文组织答案，确保内容完整、重点突出、易于理解。"""

        elif "说明" in user_input or "解释" in user_input or "启示" in user_input:
            # 为说明/解释/启示类查询
            prompt = f"""请基于以下知识库内容，深入说明和阐释{subject}对现代人的启示和应用价值：

【知识库搜索结果】
{search_result}

【说明要求】
1. 概括核心观点和主要思想
2. 分析这些观点的现代意义和应用价值
3. 说明对当代人生活、工作、思想的启示
4. 列举具体的应用场景或案例
5. 阐述这些思想与现代价值观的关联

请用深入浅出的Markdown格式组织答案，确保逻辑清晰、易于理解。"""

        else:  # 默认总结
            prompt = f"""请对以下知识库内容进行全面的总结和结构化整理：

【知识库搜索结果】
{search_result}

【整理要求】
1. 提取主要观点、论据和论证逻辑
2. 组织成清晰的层级结构（使用标题、子标题、列表等）
3. 突出重点，删除冗余内容
4. 保留重要引用和出处信息
5. 添加简要的段落摘要

请用结构化的Markdown格式输出，确保内容完整、逻辑清晰、易于理解。"""

        try:
            logger.info(f"[KB] 执行LLM后处理: {user_input}... (timeout={timeout}s)")
            # 调用LLM进行处理，使用指定的超时时间和模型
            if model:
                processed_output = self.model_caller(prompt, model=model, temperature=0.3, timeout=timeout)
            else:
                processed_output = self.model_caller(prompt, temperature=0.3, timeout=timeout)
            logger.info(f"[KB] LLM后处理完成，输出长度: {len(processed_output)}")
            return processed_output
        except Exception as e:
            logger.warning(f"[KB] LLM后处理失败: {e}，返回原始结果")
            return search_result

    def _save_and_return_download_link(self, content: str, user_input: str, output_format: str = "md", context: Optional[Dict] = None) -> str:
        """将搜索结果保存为文件并返回下载链接

        Args:
            content: 要保存的内容
            user_input: 用户原始输入（用于生成文件名）
            output_format: 输出格式 (md, docx, pptx)
            context: 上下文信息（可能包含 content_to_convert）
        """
        # os, datetime 已在文件顶部导入

        try:
            # 检查是否有需要转换的历史内容
            if context and context.get('content_to_convert'):
                history_content = context.get('content_to_convert')
                logger.info(f"[KB] 使用历史对话中的内容进行转换 (长度: {len(history_content)})")
                content = history_content
            # 清理用户输入，移除文件名中不安全的字符和可能干扰Markdown链接的字符
            safe_query = user_input.replace("/", "_").replace("\\", "_").replace(":", "_")
            safe_query = re.sub(r'[<>:"*?()]', '', safe_query)  # 移除括号和其他特殊字符
            safe_query = safe_query[:30]
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

            save_dir = "/home/will/Downloads/opencode_p/tmp/"
            os.makedirs(save_dir, exist_ok=True)

            filename = None
            file_path = None

            if output_format == "md":
                filename = f"搜索结果_{safe_query}_{timestamp}.md"
                file_path = os.path.join(save_dir, filename)
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.write(content)

            elif output_format == "docx":
                filename = f"搜索结果_{safe_query}_{timestamp}.docx"
                file_path = os.path.join(save_dir, filename)
                self._create_docx(content, file_path)

            elif output_format == "pptx":
                filename = f"搜索结果_{safe_query}_{timestamp}.pptx"
                file_path = os.path.join(save_dir, filename)
                self._create_pptx(content, file_path)

            logger.info(f"[KB] 文件已保存: {file_path}")

            encoded_filename = url_quote(str(filename))
            base_url = context.get("base_url", "") if context else ""
            download_url = f"{base_url}/api/download/{encoded_filename}"

            # 转义 Markdown 链接文本中的特殊字符
            # 圆括号在 Markdown 链接的显示文本中有特殊含义，需要转义
            display_filename = filename.replace('(', r'\(').replace(')', r'\)').replace('[', r'\[').replace(']', r'\]')

            format_names = {"md": "Markdown", "docx": "Word文档", "pptx": "PowerPoint演示文稿"}
            format_name = format_names.get(output_format, output_format.upper())

            return f"""{content}

---

📥 **下载{format_name}**
[{format_name}: {display_filename}]{download_url}
"""

        except Exception as e:
            logger.error(f"[KB] 文件保存失败: {e}")
            return f"{content}\n\n⚠️ 文件保存失败: {e}"
#   工作流程：                                                                                                                                      
#   用户请求 "总结道德经并生成Word文档"                                                                                                             
#       ↓                                                                                                                                           
#   1. 知识库搜索 → 获取相关内容                                                                                                                    
#       ↓                                                                                                                                           
#   2. LLM后处理 → 生成结构化总结                                                                                                                   
#       ↓                                                                                                                                           
#   3. 检测到需要docx格式                                                                                                                           
#       ↓                                                                                                                                           
#   4. 保存内容为临时markdown文件                                                                                                                   
#       ↓                                                                                                                                           
#   5. 调用Claude使用docx skill                                                                                                                   
#       ↓                                                                                                                                           
#   6. docx skill使用docx-js创建Word文档                                                                                                            
#       ↓                                                                                                                                           
#   7. 返回下载链接给用户     

#应该用skill 
    def _create_docx(self, content: str, file_path: str):
        """使用Claude调用docx skill创建Word文档

        通过Claude调用docx skill，而不是生成Python代码执行
        注意：将详细prompt写入临时文件，避免shell命令行转义问题
        """
        import subprocess
        import time

        start_time = time.time()

        # 先将内容保存为临时markdown文件
        temp_md_path = file_path.replace('.docx', '_temp.md')
        with open(temp_md_path, 'w', encoding='utf-8') as f:
            f.write(content)

        logger.info(f"[DOCX] 开始调用Claude使用docx skill创建文档: {file_path}")
        logger.info(f"[DOCX] 临时markdown文件: {temp_md_path}, 内容长度: {len(content)}")

        # 构建Claude调用命令
        prompt = f"""请使用docx skill将以下markdown内容创建为Word文档。

输出文件路径: {file_path}

markdown内容已保存在: {temp_md_path}

请读取该markdown文件的内容，然后使用docx-js创建Word文档，保存到指定路径。
文档应该保持markdown的结构：
- # 标题转为一级标题
- ## 标题转为二级标题
- ### 标题转为三级标题
- 列表项保持列表格式
- markdown表格更换为word表格
- **xx**转为加粗的xx
- 普通文本作为段落

创建完成后输出"文档已创建: {file_path}"
"""

        # 将prompt写入临时文件，避免shell命令行转义问题
        temp_prompt_path = file_path.replace('.docx', '_prompt.txt')
        with open(temp_prompt_path, 'w', encoding='utf-8') as f:
            f.write(prompt)

        claude_executable = os.environ.get('CLAUDE_EXECUTABLE', 'claude')

        try:
            # 使用简短命令引用prompt文件，避免shell转义问题
            simple_prompt = f"请阅读文件 {temp_prompt_path} 中的完整指令，按照指令使用docx skill创建Word文档，源内容在 {temp_md_path}，输出到 {file_path}"

            logger.info(f"[DOCX] 执行Claude命令，prompt文件: {temp_prompt_path}")
            logger.info(f"[DOCX] subprocess.run 开始，timeout=300秒")

            proc_start = time.time()
            result = subprocess.run(
                [claude_executable, '--print', '--model', 'haiku'],
                input=simple_prompt,  # 通过 stdin 传递输入
                capture_output=True,
                text=True,
                timeout=300,
                cwd="/home/will/Downloads/opencode_p"
            )
            proc_elapsed = time.time() - proc_start

            logger.info(f"[DOCX] subprocess.run 完成，耗时: {proc_elapsed:.1f}秒, returncode: {result.returncode}")
            if result.stdout:
                logger.info(f"[DOCX] stdout (前500字): {result.stdout}")
            if result.stderr:
                logger.warning(f"[DOCX] stderr: {result.stderr}")

            # 检查文件是否创建成功
            if os.path.exists(file_path):
                file_size = os.path.getsize(file_path)
                total_elapsed = time.time() - start_time
                logger.info(f"[DOCX] 文档创建成功: {file_path}, 大小: {file_size} bytes, 总耗时: {total_elapsed:.1f}秒")
            else:
                # 如果Claude没有成功创建，回退到简单的python-docx方式
                logger.warning(f"[DOCX] Claude未能创建文档（文件不存在），使用回退方案")
                self._create_docx_fallback(content, file_path)

            # 清理临时文件
            for tmp in [temp_md_path, temp_prompt_path]:
                if os.path.exists(tmp):
                    os.remove(tmp)

        except subprocess.TimeoutExpired as e:
            proc_elapsed = time.time() - start_time
            logger.warning(f"[DOCX] subprocess.run 超时! 已运行: {proc_elapsed:.1f}秒, timeout设置: 300秒")
            logger.warning(f"[DOCX] TimeoutExpired详情: {e}")
            self._create_docx_fallback(content, file_path)
            for tmp in [temp_md_path, temp_prompt_path]:
                if os.path.exists(tmp):
                    os.remove(tmp)
        except Exception as e:
            proc_elapsed = time.time() - start_time
            logger.error(f"[DOCX] Claude调用失败: {e}, 耗时: {proc_elapsed:.1f}秒，使用回退方案")
            self._create_docx_fallback(content, file_path)
            for tmp in [temp_md_path, temp_prompt_path]:
                if os.path.exists(tmp):
                    os.remove(tmp)

    def _create_docx_fallback(self, content: str, file_path: str):
        """回退方案：使用python-docx创建简单的Word文档，支持markdown表格和格式"""
        try:
            from docx import Document
            from docx.shared import Pt
            import re

            doc = Document()
            lines = content.split('\n')

            # 用于收集表格行
            table_lines = []
            in_table = False

            def is_table_line(line):
                """判断是否是markdown表格行"""
                stripped = line.strip()
                return stripped.startswith('|') and stripped.endswith('|')

            def is_separator_line(line):
                """判断是否是表格分隔行 (如 |---|---|)"""
                stripped = line.strip()
                if not stripped.startswith('|') or not stripped.endswith('|'):
                    return False
                # 检查是否只包含 |, -, :, 空格
                inner = stripped[1:-1]
                return all(c in '-|: ' for c in inner)

            def add_formatted_text(paragraph, text):
                """解析markdown格式并添加到段落中，支持粗体、斜体、代码等"""
                # 匹配模式：**粗体**、*斜体*、`代码`、~~删除线~~
                # 使用正则表达式分割文本，保留匹配的格式标记
                pattern = r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`|~~[^~]+~~)'
                parts = re.split(pattern, text)

                for part in parts:
                    if not part:
                        continue

                    if part.startswith('**') and part.endswith('**'):
                        # 粗体
                        run = paragraph.add_run(part[2:-2])
                        run.bold = True
                    elif part.startswith('*') and part.endswith('*') and not part.startswith('**'):
                        # 斜体
                        run = paragraph.add_run(part[1:-1])
                        run.italic = True
                    elif part.startswith('`') and part.endswith('`'):
                        # 代码（使用等宽字体）
                        run = paragraph.add_run(part[1:-1])
                        run.font.name = 'Consolas'
                    elif part.startswith('~~') and part.endswith('~~'):
                        # 删除线
                        run = paragraph.add_run(part[2:-2])
                        run.font.strike = True
                    else:
                        # 普通文本
                        paragraph.add_run(part)

            def add_table_to_doc(doc, table_lines):
                """将收集的表格行转换为Word表格"""
                if not table_lines:
                    return

                # 解析表格数据
                rows_data = []
                for line in table_lines:
                    if is_separator_line(line):
                        continue  # 跳过分隔行
                    # 解析单元格
                    cells = [cell.strip() for cell in line.strip('|').split('|')]
                    rows_data.append(cells)

                if not rows_data:
                    return

                # 确定列数
                num_cols = max(len(row) for row in rows_data)
                num_rows = len(rows_data)

                # 创建表格
                table = doc.add_table(rows=num_rows, cols=num_cols)
                table.style = 'Table Grid'

                # 填充表格
                for i, row_data in enumerate(rows_data):
                    row = table.rows[i]
                    for j, cell_text in enumerate(row_data):
                        if j < num_cols:
                            cell = row.cells[j]
                            # 清空默认段落并添加格式化文本
                            cell.paragraphs[0].clear()
                            add_formatted_text(cell.paragraphs[0], cell_text)
                            # 第一行加粗（表头）
                            if i == 0:
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.bold = True

                doc.add_paragraph()  # 表格后添加空行

            i = 0
            while i < len(lines):
                line = lines[i]
                line_stripped = line.strip()

                # 检查是否是表格行
                if is_table_line(line_stripped):
                    if not in_table:
                        in_table = True
                        table_lines = []
                    table_lines.append(line_stripped)
                    i += 1
                    continue
                else:
                    # 如果之前在表格中，现在不是表格行了，处理表格
                    if in_table:
                        add_table_to_doc(doc, table_lines)
                        table_lines = []
                        in_table = False

                # 处理普通行
                if not line_stripped:
                    doc.add_paragraph()
                elif line_stripped.startswith('# '):
                    heading = doc.add_heading(level=1)
                    add_formatted_text(heading, line_stripped[2:])
                elif line_stripped.startswith('## '):
                    heading = doc.add_heading(level=2)
                    add_formatted_text(heading, line_stripped[3:])
                elif line_stripped.startswith('### '):
                    heading = doc.add_heading(level=3)
                    add_formatted_text(heading, line_stripped[4:])
                elif line_stripped.startswith('#### '):
                    heading = doc.add_heading(level=4)
                    add_formatted_text(heading, line_stripped[5:])
                elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
                    p = doc.add_paragraph(style='List Bullet')
                    add_formatted_text(p, line_stripped[2:])
                elif len(line_stripped) > 2 and line_stripped[0].isdigit() and line_stripped[1] == '.':
                    p = doc.add_paragraph(style='List Number')
                    add_formatted_text(p, line_stripped[2:].strip())
                elif line_stripped.startswith('> '):
                    # 引用块
                    p = doc.add_paragraph()
                    p.style = 'Quote' if 'Quote' in [s.name for s in doc.styles] else 'Normal'
                    add_formatted_text(p, line_stripped[2:])
                elif line_stripped.startswith('---') or line_stripped.startswith('***'):
                    # 分隔线 - 添加一个空段落
                    doc.add_paragraph('─' * 50)
                else:
                    # 普通段落 - 支持格式化
                    p = doc.add_paragraph()
                    add_formatted_text(p, line_stripped)

                i += 1

            # 处理文档末尾可能残留的表格
            if in_table and table_lines:
                add_table_to_doc(doc, table_lines)

            doc.save(file_path)
            logger.info(f"[DOCX] ✅ 回退方案创建文档成功: {file_path}")
        except ImportError:
            logger.error("[DOCX] python-docx未安装，无法创建文档")
            raise Exception("python-docx库未安装，请运行: pip install python-docx")
        except Exception as e:
            logger.error(f"[DOCX] 回退方案失败: {e}")
            raise e

    def _create_pptx(self, content: str, file_path: str):
        """使用Claude调用pptx skill创建PowerPoint演示文稿

        遵循 pptx/SKILL.md 的工作流程：
        1. 分析内容结构，设计大纲
        2. 选择配色方案和布局
        3. 创建HTML幻灯片
        4. 使用html2pptx转换为PPTX

        注意：将详细prompt写入临时文件，避免shell命令行转义问题
        """
        import subprocess
        import time

        start_time = time.time()

        # 先将内容保存为临时markdown文件
        temp_md_path = file_path.replace('.pptx', '_temp.md')
        with open(temp_md_path, 'w', encoding='utf-8') as f:
            f.write(content)

        # 分析内容结构，生成大纲
        lines = content.split('\n')
        sections = []
        current_section = {"title": "", "content": []}

        for line in lines:
            if line.startswith('# '):
                if current_section["title"] or current_section["content"]:
                    sections.append(current_section)
                current_section = {"title": line[2:].strip(), "content": [], "level": 1}
            elif line.startswith('## '):
                if current_section["title"] or current_section["content"]:
                    sections.append(current_section)
                current_section = {"title": line[3:].strip(), "content": [], "level": 2}
            elif line.startswith('### '):
                current_section["content"].append({"type": "subtitle", "text": line[4:].strip()})
            elif line.startswith('- ') or line.startswith('* '):
                current_section["content"].append({"type": "bullet", "text": line[2:].strip()})
            elif line.startswith('**') and line.endswith('**'):
                current_section["content"].append({"type": "emphasis", "text": line[2:-2].strip()})
            elif line.strip():
                current_section["content"].append({"type": "text", "text": line.strip()})

        if current_section["title"] or current_section["content"]:
            sections.append(current_section)

        # 生成大纲描述
        outline_text = "内容大纲：\n"
        for i, sec in enumerate(sections[:10]):  # 最多10个章节
            outline_text += f"  幻灯片{i+1}: {sec.get('title', '无标题')}\n"
            for item in sec.get('content', [])[:5]:  # 每章节最多5个要点
                outline_text += f"    - {item['text']}...\n" if len(item['text']) > 50 else f"    - {item['text']}\n"

        logger.info(f"[PPTX] 开始调用Claude使用pptx skill创建演示文稿: {file_path}")
        logger.info(f"[PPTX] 临时markdown文件: {temp_md_path}, 内容长度: {len(content)}, 识别到 {len(sections)} 个章节")

        # 构建详细的Claude调用提示词，遵循SKILL.md工作流程
        prompt = f"""请使用pptx skill将以下markdown内容创建为专业的PowerPoint演示文稿。

## 任务说明
输出文件路径: {file_path}
源文件路径: {temp_md_path}

## 工作流程（必须遵循）

### 步骤1：阅读skill指南
首先阅读 /home/will/.claude/skills/pptx/html2pptx.md 了解HTML幻灯片语法。

### 步骤2：内容分析与大纲设计
{outline_text}

### 步骤3：选择设计方案
根据内容主题选择合适的配色方案（从SKILL.md中的16个配色方案中选择）：
- 商业/企业主题：Classic Blue 或 Charcoal and Red
- 科技/创新主题：Deep Purple and Emerald 或 Vibrant Orange
- 专业/严肃主题：Black and Gold 或 Forest Green
- 温暖/亲和主题：Warm Blush 或 Cream and Forest Green

### 步骤4：创建HTML幻灯片
为每个章节创建HTML幻灯片，遵循以下规则：
- 幻灯片尺寸: 720pt x 405pt (16:9)
- 所有文字必须在 p, h1-h6, ul, ol 标签内
- 只使用web-safe字体: Arial, Verdana, Georgia, Times New Roman
- 使用flexbox布局
- 一级标题 转为 标题幻灯片（大字居中）
- 二级标题 转为 新幻灯片的标题
- 三级标题 转为 幻灯片内的子标题
- 列表项用 ul/li 格式，不要手动添加符号
- 加粗文字用 b 或 strong 标签

### 步骤5：使用html2pptx.js转换
创建JavaScript文件，使用html2pptx库将HTML幻灯片转换为PowerPoint并保存。

### 步骤6：验证
确认文件 {file_path} 已创建。

## 重要提示
- 保持专业的视觉效果
- 每张幻灯片内容不要过多（3-5个要点）
- 使用一致的字体大小和颜色
- 创建完成后输出"演示文稿已创建: {file_path}"
"""

        # 将prompt写入临时文件，避免shell命令行转义问题
        temp_prompt_path = file_path.replace('.pptx', '_prompt.txt')
        with open(temp_prompt_path, 'w', encoding='utf-8') as f:
            f.write(prompt)

        claude_executable = os.environ.get('CLAUDE_EXECUTABLE', 'claude')

        try:
            # 使用简短命令引用prompt文件，避免shell转义问题
            simple_prompt = f"请阅读文件 {temp_prompt_path} 中的完整指令，按照指令使用pptx skill创建PowerPoint演示文稿，源内容在 {temp_md_path}，输出到 {file_path}"

            logger.info(f"[PPTX] 执行Claude命令，prompt文件: {temp_prompt_path}")
            logger.info(f"[PPTX] subprocess.run 开始，timeout=300秒")

            proc_start = time.time()
            result = subprocess.run(
                [claude_executable, '--print', '--model', 'haiku'],
                input=simple_prompt,  # 通过 stdin 传递输入
                capture_output=True,
                text=True,
                timeout=300,
                cwd="/home/will/Downloads/opencode_p"
            )
            proc_elapsed = time.time() - proc_start

            logger.info(f"[PPTX] subprocess.run 完成，耗时: {proc_elapsed:.1f}秒, returncode: {result.returncode}")
            if result.stdout:
                logger.info(f"[PPTX] stdout (前500字): {result.stdout}")
            if result.stderr:
                logger.warning(f"[PPTX] stderr: {result.stderr}")

            # 检查文件是否创建成功
            if os.path.exists(file_path):
                file_size = os.path.getsize(file_path)
                total_elapsed = time.time() - start_time
                logger.info(f"[PPTX] 演示文稿创建成功: {file_path}, 大小: {file_size} bytes, 总耗时: {total_elapsed:.1f}秒")
            else:
                # 如果Claude没有成功创建，回退到简单方式
                logger.warning(f"[PPTX] Claude未能创建演示文稿（文件不存在），使用回退方案")
                self._create_pptx_fallback(content, file_path)

            # 清理临时文件
            for tmp in [temp_md_path, temp_prompt_path]:
                if os.path.exists(tmp):
                    os.remove(tmp)

        except subprocess.TimeoutExpired as e:
            proc_elapsed = time.time() - start_time
            logger.warning(f"[PPTX] subprocess.run 超时! 已运行: {proc_elapsed:.1f}秒, timeout设置: 300秒")
            logger.warning(f"[PPTX] TimeoutExpired详情: {e}")
            self._create_pptx_fallback(content, file_path)
            for tmp in [temp_md_path, temp_prompt_path]:
                if os.path.exists(tmp):
                    os.remove(tmp)
        except Exception as e:
            proc_elapsed = time.time() - start_time
            logger.error(f"[PPTX] Claude调用失败: {e}, 耗时: {proc_elapsed:.1f}秒，使用回退方案")
            self._create_pptx_fallback(content, file_path)
            for tmp in [temp_md_path, temp_prompt_path]:
                if os.path.exists(tmp):
                    os.remove(tmp)

    def _create_pptx_fallback(self, content: str, file_path: str):
        """回退方案：使用python-pptx创建简单的PowerPoint演示文稿，支持markdown格式"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            import re

            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9 宽度
            prs.slide_height = Inches(7.5)    # 16:9 高度

            lines = content.split('\n')
            current_content = []

            def add_formatted_text_to_paragraph(paragraph, text):
                """解析markdown格式并添加到段落中，支持粗体、斜体、代码等"""
                # 匹配模式：**粗体**、*斜体*、`代码`、~~删除线~~
                pattern = r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`|~~[^~]+~~)'
                parts = re.split(pattern, text)

                for part in parts:
                    if not part:
                        continue

                    if part.startswith('**') and part.endswith('**'):
                        # 粗体
                        run = paragraph.add_run(part[2:-2])
                        run.bold = True
                    elif part.startswith('*') and part.endswith('*') and not part.startswith('**'):
                        # 斜体
                        run = paragraph.add_run(part[1:-1])
                        run.italic = True
                    elif part.startswith('`') and part.endswith('`'):
                        # 代码（使用等宽字体，如果可用）
                        run = paragraph.add_run(part[1:-2])
                        try:
                            run.font.name = 'Consolas'
                        except:
                            pass
                    elif part.startswith('~~') and part.endswith('~~'):
                        # 删除线
                        run = paragraph.add_run(part[2:-2])
                        # python-pptx 不直接支持删除线，跳过
                    else:
                        # 普通文本
                        paragraph.add_run(part)

            def add_slide_with_content(prs, title, content_list):
                """添加带内容的幻灯片"""
                # 使用标题和内容布局
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                # 设置标题
                if title:
                    slide.shapes.title.text = title
                    # 标题使用大字号
                    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)

                # 添加内容
                if content_list and slide.placeholders[1]:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()  # 清空默认内容

                    for i, item in enumerate(content_list):
                        if i == 0:
                            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                        else:
                            p = tf.add_paragraph()

                        p.text = ""  # 清空默认文本
                        add_formatted_text_to_paragraph(p, item)
                        p.level = 0

                        # 设置字体大小
                        for run in p.runs:
                            run.font.size = Pt(18)

                return slide

            def add_title_slide(prs, title, subtitle=""):
                """添加标题幻灯片"""
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)

                # 添加标题
                if title:
                    left = Inches(0.5)
                    top = Inches(2)
                    width = Inches(12.333)
                    height = Inches(1.5)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    add_formatted_text_to_paragraph(p, title)
                    for run in p.runs:
                        run.font.size = Pt(44)
                        run.font.bold = True

                # 添加副标题
                if subtitle:
                    left = Inches(0.5)
                    top = Inches(4)
                    width = Inches(12.333)
                    height = Inches(1)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    add_formatted_text_to_paragraph(p, subtitle)
                    for run in p.runs:
                        run.font.size = Pt(24)

                return slide

            current_title = ""
            first_slide = True

            for line in lines:
                line_stripped = line.strip()

                if not line_stripped:
                    continue

                if line_stripped.startswith('# '):
                    # 一级标题 - 标题幻灯片
                    if current_title or current_content:
                        if first_slide:
                            add_title_slide(prs, current_title)
                            first_slide = False
                        else:
                            add_slide_with_content(prs, current_title, current_content)
                        current_content = []
                    current_title = line_stripped[2:]
                elif line_stripped.startswith('## ') or line_stripped.startswith('### '):
                    # 二级/三级标题 - 新幻灯片
                    if current_title or current_content:
                        if first_slide:
                            add_title_slide(prs, current_title)
                            first_slide = False
                        else:
                            add_slide_with_content(prs, current_title, current_content)
                        current_content = []
                    current_title = line_stripped.lstrip('#').strip()
                elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
                    current_content.append(line_stripped[2:])
                elif len(line_stripped) > 2 and line_stripped[0].isdigit() and line_stripped[1] == '.':
                    current_content.append(line_stripped[2:].strip())
                else:
                    current_content.append(line_stripped)

            # 添加最后一个幻灯片
            if current_title or current_content:
                if first_slide:
                    add_title_slide(prs, current_title)
                else:
                    add_slide_with_content(prs, current_title, current_content)

            # 如果没有任何幻灯片，添加一个默认的
            if len(prs.slides) == 0:
                slide_layout = prs.slide_layouts[5]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(12)
                height = Inches(6)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                add_formatted_text_to_paragraph(tf.paragraphs[0], content)  # 限制内容长度

            prs.save(file_path)
            logger.info(f"[PPTX] ✅ 回退方案创建演示文稿成功: {file_path}, 共 {len(prs.slides)} 张幻灯片")
        except ImportError:
            logger.error("[PPTX] python-pptx未安装，无法创建演示文稿")
            raise Exception("python-pptx库未安装，请运行: pip install python-pptx")
        except Exception as e:
            logger.error(f"[PPTX] 回退方案失败: {e}")
            raise e

    def _dispatch_web_search_dk(self, user_input: str, context: Optional[Dict]) -> str:
        """使用 dk-search 服务进行网络搜索 - 自动优化查询词"""
        import httpx
        from datetime import datetime
        import logging
        logger = logging.getLogger(__name__)

        logger.info(f"[WebSearch-DK] 处理网络搜索请求: {user_input[:80]}...")

        # ========== 1. 查询词优化 ==========
        query = user_input.strip()

        # 1.1 检查是否包含日期/时间相关词汇
        date_keywords = [
            '今天', '今日', '昨天', '前天', '本周', '上周', '本月', '上月',
            '最近', '最新', '刚刚', '近期', '目前', '现在',
            '2024', '2025', '2026', '2027', '年', '月', '日',
            'today', 'yesterday', 'this week', 'this month', 'recent', 'latest', 'now', 'current'
        ]
        has_date = any(keyword in query.lower() for keyword in date_keywords)

        # 1.2 如果没有日期信息，自动添加当前年月
        if not has_date:
            current_date = datetime.now().strftime("%Y年%m月")
            query = f"{query} {current_date}"
            logger.info(f"[WebSearch-DK] 自动添加日期，查询变更为: '{query}'")

        # 1.3 提取搜索关键词（移除口语化表达）
        remove_patterns = [
            '请帮我', '帮我', '请', '我想', '我要', '能否', '可以',
            '搜索一下', '查找一下', '找一下', '查一下', '搜一下',
            '从网络搜索', '网络搜索', '在网上搜', '上网查',
            'please', 'help me', 'can you', 'could you'
        ]
        optimized_query = query
        for pattern in remove_patterns:
            optimized_query = optimized_query.replace(pattern, '')
        optimized_query = ' '.join(optimized_query.split())  # 清理多余空格

        if optimized_query != query:
            logger.info(f"[WebSearch-DK] 优化查询词: '{query}' -> '{optimized_query}'")
            query = optimized_query
        logger.info(f"优化后的查询词：{query}")
        # ========== 2. 调用 dk-search 服务 ==========
        dk_search_url = "http://localhost:8001/fetch"
        provider = "ddg"  # 默认使用 DuckDuckGo
        max_results = 5

        try:
            with httpx.Client(timeout=60) as client:
                logger.info(f"[WebSearch-DK] 调用 dk-search: query='{query}', provider={provider}, max_results={max_results}")
                response = client.post(
                    dk_search_url,
                    json={
                        "query": query,
                        "provider": provider,
                        "max_results": max_results,
                        "timeout": 15
                    }
                )

                if response.status_code != 200:
                    logger.error(f"[WebSearch-DK] dk-search 返回错误: {response.status_code}")
                    return f"⚠️ 网络搜索服务返回错误 (状态码: {response.status_code})\n\n建议您稍后重试或直接访问搜索引擎。"

                result = response.json()
                results = result.get("results", [])
                logger.info(f"DK返回的结果:{results[:500]}")

                if not results:
                    logger.warning(f"[WebSearch-DK] 未找到搜索结果")
                    return f"🔍 搜索完成，但未找到相关结果。\n\n查询: {query}\n\n建议尝试不同的关键词或直接访问相关网站。"

                # ========== 3. 格式化搜索结果 ==========
                formatted_output = []
                formatted_output.append(f"🔍 **搜索结果** (查询: {query})\n")
                formatted_output.append(f"共找到 {len(results)} 条相关信息：\n")
                formatted_output.append("-" * 50)

                for i, r in enumerate(results, 1):
                    title = r.get("title", "无标题")
                    url = r.get("url", "")
                    content = r.get("content", "")
                    snippet = r.get("snippet", "")

                    formatted_output.append(f"\n### [{i}] {title}")
                    if url:
                        formatted_output.append(f"🔗 {url}")

                    # 优先使用内容，其次使用摘要
                    text = content if content else snippet
                    if text:
                        # 截取前800字符作为预览
                        preview = text[:800] + "..." if len(text) > 800 else text
                        formatted_output.append(f"\n{preview}")

                    formatted_output.append("")

                formatted_output.append("-" * 50)
                formatted_output.append(f"\n📅 搜索时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
                formatted_output.append(f"🔎 搜索引擎: {provider.upper()}")

                final_result = "\n".join(formatted_output)
                logger.info(f"[WebSearch-DK] ✅ 搜索完成，获取 {len(results)} 条结果，返回 {len(final_result)} 字符\n{final_result}")
                prompt = f"""请根据{user_input}对下面搜索到的网页结果按主题、重要性和时间（最新的在最前）进行总结：
                {final_result}
                要求：使用markdown形式，如有表格以表格呈现；不得添加网页结果中没有的内容。
                """
                model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"
                return self.model_caller(prompt, model=model, temperature=0.3)


        except httpx.ConnectError:
            logger.error(f"[WebSearch-DK] ❌ 无法连接到 dk-search 服务 (localhost:8001)")
            return "⚠️ 网络搜索服务不可用\n\ndk-search 服务未启动，请确保服务正在运行 (端口 8001)。\n\n您可以尝试：\n1. 启动 dk-search 服务\n2. 使用其他搜索方式"
        except Exception as e:
            logger.error(f"[WebSearch-DK] ❌ 搜索失败: {e}", exc_info=True)
            return f"⚠️ 网络搜索失败: {str(e)}\n\n建议您直接访问搜索引擎查询最新信息。"

    def _dispatch_web_fetch(self, user_input: str, context: Optional[Dict]) -> str:
        """获取网页内容（如 GitHub URL）- 使用 httpx 抓取并用 LLM 整理"""
        import httpx
        import re
        import logging
        logger = logging.getLogger(__name__)

        logger.info(f"[WebFetch] 处理网页获取请求: {user_input[:80]}...")

        # 1. 从 user_input 中提取 URL
        url_pattern = r'https?://[^\s<>"{}|\\^`\[\]]+'
        urls = re.findall(url_pattern, user_input)

        if not urls:
            logger.warning(f"[WebFetch] 未找到 URL，降级为 Claude 直接处理")
            return self._dispatch_direct(user_input, context)

        url = urls[0]
        logger.info(f"[WebFetch] 获取 URL: {url}")

        # 2. 对 GitHub URL 进行优化，尝试获取 raw 内容
        fetch_urls = [url]
        if "github.com" in url and "/blob/" in url:
            # 将 GitHub blob URL 转为 raw URL
            raw_url = url.replace("github.com", "raw.githubusercontent.com").replace("/blob/", "/")
            fetch_urls.insert(0, raw_url)
        elif "github.com" in url and not any(x in url for x in ["/blob/", "/raw/", "/tree/"]):
            # 仓库根目录，尝试获取 README
            readme_url = url.rstrip("/") + "/raw/main/README.md"
            fetch_urls.insert(0, readme_url)

        # 3. 依次尝试获取内容
        content = None
        fetched_url = None
        try:
            with httpx.Client(timeout=30, follow_redirects=True) as client:
                for try_url in fetch_urls:
                    try:
                        logger.info(f"[WebFetch] 尝试获取: {try_url}")
                        response = client.get(try_url, headers={
                            "User-Agent": "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36",
                            "Accept": "text/html,application/xhtml+xml,text/plain,text/markdown,*/*"
                        })
                        if response.status_code == 200:
                            content = response.text
                            fetched_url = try_url
                            break
                        else:
                            logger.debug(f"[WebFetch] {try_url} 返回 {response.status_code}")
                    except Exception as e:
                        logger.debug(f"[WebFetch] {try_url} 获取失败: {e}")
                        continue
        except Exception as e:
            logger.error(f"[WebFetch] 网页获取失败: {e}")
            return f"⚠️ 无法获取网页内容: {str(e)}\n\nURL: {url}"

        if not content:
            return f"⚠️ 无法获取网页内容（所有尝试均失败）\n\nURL: {url}"

        # 4. 截取内容（避免过长）
        max_len = 8000
        if len(content) > max_len:
            content = content[:max_len] + f"\n\n... (内容截断，原始长度 {len(content)} 字符)"

        logger.info(f"[WebFetch] ✅ 获取成功: {fetched_url}, 内容长度 {len(content)} 字符")

        # 5. 使用 LLM 整理内容
        prompt = f"""请根据用户的需求，整理和提取以下网页内容中的关键信息。

用户请求: {user_input}

网页来源: {fetched_url}
网页内容:
{content}

要求：
- 提取与用户需求相关的核心内容
- 保持代码块的完整性
- 使用 markdown 格式组织输出
- 标明信息来源 URL"""

        model = context.get("model", "nvidia/minimaxai/minimax-m2.1") if context else "nvidia/minimaxai/minimax-m2.1"
        return self.model_caller(prompt, model=model, temperature=0.3)

    #Claude内置网络搜索保留但不用
    def _dispatch_web_search(self, user_input: str, context: Optional[Dict]) -> str:
        """处理网络搜索请求 - 调用 Claude 使用 WebSearch 工具获取实时信息"""
        import logging
        logger = logging.getLogger(__name__)

        logger.info(f"[WebSearch] 处理网络搜索请求: {user_input[:80]}...")

        # 构建网络搜索提示，明确指示使用 WebSearch 工具获取实时信息
        search_prompt = f"""用户需要获取实时信息：{user_input}

请通过以下步骤完成任务：
1. 理解用户的查询需求
2. 使用 WebSearch 工具搜索相关信息
3. 从搜索结果中提取关键信息
4. 用清晰、结构化的方式总结和呈现信息，包括：
   - 最新的相关数据
   - 关键时间戳或时间范围
   - 信息来源
   - 必要的背景说明

请使用您的 WebSearch 工具来获取最新的网络信息。"""

        # 从context中获取用户选择的模型，如果没有则使用默认值
        model = "haiku" #context.get("model", "haiku") if context else 

        try:
            # 调用 model_caller，网络搜索需要更长超时（10分钟）
            result = self.model_caller(search_prompt, model=model, temperature=0.3, timeout=600)
            logger.info(f"[WebSearch] 网络搜索完成，返回 {len(result)} 字内容")
            return result
        except Exception as e:
            logger.error(f"[WebSearch] ❌ 网络搜索失败: {e}")
            # 降级处理：返回错误信息
            return f"⚠️ 网络搜索失败: {str(e)}\n\n建议您直接访问相关网站查询最新信息。"

    def _dispatch_direct(self, user_input: str, context: Optional[Dict]) -> str:
        """直接让 Claude 处理"""
        # 从context中获取用户选择的模型，如果没有则使用默认值
        logger.info(f"==该步direct执行:{str(user_input)}")
        model = "haiku" #context.get("model", "haiku") if context else
        
        # 构建完整的 prompt，包含对话历史
        full_prompt = self._build_prompt_with_history(user_input, context)
        result0 = self.model_caller(full_prompt, model=model, temperature=0.5, timeout=300)
        logger.info(f"===========dispatch 结果：{result0}")
        result = self._execute_docx_code(result0,context)
        logger.info(f"===========dispatch 执行结果：{result}")
        return result

    def _build_prompt_with_history(self, user_input: str, context: Optional[Dict]) -> str:
        """构建包含对话历史的完整 prompt"""
        if not context:
            return user_input

        history = context.get('history', [])
        if not history:
            return user_input

        # 构建历史对话文本
        history_parts = []
        for msg in history[-10:]:  # 最多取最近10条历史
            role = msg.get('role', 'user')
            content = msg.get('content', '')
            if content:
                if role == 'user':
                    history_parts.append(f"用户: {content}")
                elif role == 'assistant':
                    # 对于 assistant 内容，如果太长则截取
                    logger.info(f"assistant:{content}")
                    if len(content) > 3000:
                        content = content[:3000] + "...(内容过长已截断)"
                    history_parts.append(f"助手: {content}")
                elif role == 'system':
                    # 系统消息可能包含搜索结果等重要上下文
                    if len(content) > 3000:
                        content = content[:3000] + "...(内容过长已截断)"
                    history_parts.append(f"[系统信息]: {content}")

        if not history_parts:
            return user_input

        # 组合历史和当前输入
        history_text = "\n\n".join(history_parts)
        full_prompt = f"""以下是之前的对话历史，请参考这些上下文来回答用户的问题。

--- 对话历史 ---
{history_text}
--- 历史结束 ---

用户当前请求: {user_input}

请根据上述对话历史和当前请求，给出回答。"""

        logger.info(f"[Direct] 构建了包含 {len(history)} 条历史的 prompt，总长度 {len(full_prompt)} 字符")
        return full_prompt

# ============ 主入口 ============
#SemanticScheduler-SemanticAnalyzer-SkillDiscovery
#SemanticScheduler-Dispatcher

class SemanticScheduler:
    """语义调度系统 - 主入口"""

    def __init__(self, skills_dir: Path, model_caller, kb_paths: Dict = None):
        self.discovery = SkillDiscovery(skills_dir)
        self.discovery.discover_skills()

        self.analyzer = SemanticAnalyzer(self.discovery, self._wrap_model_caller(model_caller))
        kb_paths_dict = kb_paths if kb_paths else {}
        self.dispatcher = Dispatcher(self.discovery, model_caller, kb_paths_dict)

    def _wrap_model_caller(self, caller):
        """包装模型调用器以返回函数"""
        def wrapped(system_prompt: str, user_prompt: str, **kwargs):
            full_prompt = f"{system_prompt}\n\n{user_prompt}"
            return caller(full_prompt, **kwargs)
        return wrapped

    def process(self, user_input: str, context: Optional[Dict] = None) -> Dict[str, Any]:
        """处理用户输入"""
        start_time = datetime.now()

        # 获取对话历史
        history = context.get('history', []) if context else []

        # 将历史信息放入context，让dispatcher可以使用
        if context is None:
            context = {}
        context['history'] = history
        context['original_input'] = user_input  # 保存原始用户输入

        # 如果有对话历史，记录日志
        if history:
            logger.info(f"[Scheduler] 附加 {len(history)} 条对话历史")

        # 分析意图（使用原始用户输入）
        intent = self.analyzer.analyze(user_input, context)

        # 调度执行（传递原始用户输入和context，让dispatcher内部处理历史）
        result = self.dispatcher.dispatch(intent, user_input, context)

        logger.info(f"[Scheduler] 处理完成: {intent.intent_type}:{intent.target},结果：{result if result else 'None'}...")
        elapsed = (datetime.now() - start_time).total_seconds()

        return {
            "response": result,
            "intent": {
                "type": intent.intent_type,
                "target": intent.target,
                "operation": intent.operation,
                "reasoning": intent.reasoning
            },
            "metadata": {
                "processing_time": elapsed,
                "confidence": intent.confidence,
                "available_skills": len(self.discovery.skills),
                "available_agents": len(self.discovery.agents),
                "history_used": len(history) if history else 0
            }
        }

    def _format_history(self, history: List[Dict]) -> str:
        """格式化对话历史为文本"""
        if not history:
            return ""

        lines = []
        for msg in history:
            role = msg.get('role', 'unknown')
            content = msg.get('content', '')
            # 截断过长的内容
            if len(content) > 1600:
                content = content[:1600] + "..."
            role_label = "用户" if role == 'user' else "助手"
            lines.append(f"[{role_label}]: {content}")

        return "\n".join(lines)
